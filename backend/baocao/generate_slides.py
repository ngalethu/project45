#!/usr/bin/env python3
"""
Generate thesis defense PowerPoint slides for:
"Ứng dụng kiến trúc Hybrid (Edge-Cloud) và Deep Learning trong bài toán
nhận diện hành vi tài xế xe khách"
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Color palette ──
DARK_BLUE = RGBColor(0x0A, 0x2E, 0x5C)      # #0A2E5C
MEDIUM_BLUE = RGBColor(0x14, 0x4D, 0x8D)     # #144D8D
LIGHT_BLUE = RGBColor(0x1E, 0x88, 0xE5)      # #1E88E5
ACCENT_BLUE = RGBColor(0x42, 0xA5, 0xF5)     # #42A5F5
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
GRAY = RGBColor(0x9E, 0x9E, 0x9E)
DARK_GRAY = RGBColor(0x42, 0x42, 0x42)
BLACK = RGBColor(0x21, 0x21, 0x21)
ORANGE = RGBColor(0xFF, 0x6F, 0x00)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC6, 0x28, 0x28)
TEAL = RGBColor(0x00, 0x69, 0x5C)
PURPLE = RGBColor(0x6A, 0x1B, 0x9A)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


# ── Helper functions ──

def add_blank_slide():
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.fill.solid()
        if line_width:
            shape.line.width = line_width
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.fill.solid()
    return shape


def add_textbox(slide, left, top, width, height, text="", font_size=18,
                font_color=BLACK, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Calibri", line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines, font_size=16,
                          font_color=BLACK, bold=False, alignment=PP_ALIGN.LEFT,
                          font_name="Calibri", line_spacing=1.3, bullet=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(4)
        p.line_spacing = Pt(font_size * line_spacing)
        if bullet:
            p.level = 0
    return txBox


def add_slide_header(slide, title, subtitle=None):
    """Add dark blue header bar with title."""
    add_shape(slide, 0, 0, SLIDE_WIDTH, Inches(1.2), fill_color=DARK_BLUE)
    add_textbox(slide, Inches(0.8), Inches(0.15), Inches(11), Inches(0.7),
                title, font_size=32, font_color=WHITE, bold=True,
                font_name="Calibri", alignment=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(0.75), Inches(11), Inches(0.4),
                    subtitle, font_size=16, font_color=ACCENT_BLUE, bold=False,
                    font_name="Calibri", alignment=PP_ALIGN.LEFT)
    # Accent line
    add_shape(slide, 0, Inches(1.2), SLIDE_WIDTH, Inches(0.05), fill_color=LIGHT_BLUE)
    # Footer
    add_shape(slide, 0, Inches(7.1), SLIDE_WIDTH, Inches(0.4), fill_color=DARK_BLUE)
    add_textbox(slide, Inches(0.5), Inches(7.12), Inches(12), Inches(0.35),
                "La Đại Lộc | MSSV: 4551050116 | GVHD: TS. Nguyễn Thanh Bình | Trường Đại học Quy Nhơn",
                font_size=10, font_color=RGBColor(0xBB, 0xDE, 0xFB), alignment=PP_ALIGN.CENTER)


def add_bullet_items(slide, left, top, width, height, items, font_size=16,
                     font_color=DARK_GRAY, bullet_char="▸", spacing=1.4):
    """Add bullet list with custom bullet character."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{bullet_char} {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        p.line_spacing = Pt(font_size * spacing)
    return txBox


def add_card(slide, left, top, width, height, title, items, icon_text="",
             card_color=LIGHT_GRAY, title_color=DARK_BLUE, icon_bg=LIGHT_BLUE):
    """Add a card with icon, title, and bullet items."""
    card = add_rounded_rect(slide, left, top, width, height, fill_color=card_color,
                            line_color=RGBColor(0xE0, 0xE0, 0xE0))
    # Icon circle
    if icon_text:
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.2),
                                        top + Inches(0.2), Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = icon_bg
        circle.line.fill.background()
        tf = circle.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = icon_text
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)

    title_left = left + (Inches(0.9) if icon_text else Inches(0.2))
    add_textbox(slide, title_left, top + Inches(0.2), width - Inches(1), Inches(0.4),
                title, font_size=16, font_color=title_color, bold=True)

    if items:
        add_bullet_items(slide, left + Inches(0.3), top + Inches(0.7),
                         width - Inches(0.6), height - Inches(0.9),
                         items, font_size=13, bullet_char="•", spacing=1.3)


def add_speaker_notes(slide, text):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def add_arrow(slide, left, top, width, height, fill_color=LIGHT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_down_arrow(slide, left, top, width, height, fill_color=LIGHT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_flow_box(slide, left, top, width, height, text, fill_color=MEDIUM_BLUE,
                 font_color=WHITE, font_size=13):
    shape = add_rounded_rect(slide, left, top, width, height, fill_color=fill_color)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return shape


def add_table(slide, left, top, width, height, rows, cols, data, col_widths=None):
    """Add a table with data."""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c])
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.name = "Calibri"
                paragraph.alignment = PP_ALIGN.CENTER
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = DARK_GRAY
            # Header row styling
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_BLUE
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
    return table_shape


# ════════════════════════════════════════════════════════════════════
# SLIDE 1: Trang bìa
# ════════════════════════════════════════════════════════════════════
slide = add_blank_slide()

# Full background
add_shape(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, fill_color=DARK_BLUE)

# Decorative top bar
add_shape(slide, 0, 0, SLIDE_WIDTH, Inches(0.08), fill_color=LIGHT_BLUE)

# Left accent stripe
add_shape(slide, 0, 0, Inches(0.5), SLIDE_HEIGHT, fill_color=MEDIUM_BLUE)

# Title block
add_textbox(slide, Inches(1.5), Inches(0.8), Inches(10), Inches(0.6),
            "BỘ GIÁO DỤC VÀ ĐÀO TẠO", font_size=18, font_color=ACCENT_BLUE,
            bold=False, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(1.3), Inches(10), Inches(0.6),
            "TRƯỜNG ĐẠI HỌC QUY NHƠN", font_size=22, font_color=WHITE,
            bold=True, alignment=PP_ALIGN.CENTER)

# Separator line
add_shape(slide, Inches(4), Inches(2.0), Inches(5), Inches(0.04), fill_color=LIGHT_BLUE)

# Main title
add_textbox(slide, Inches(1.2), Inches(2.3), Inches(10.5), Inches(0.6),
            "KHÓA LUẬN TỐT NGHIỆP ĐẠI HỌC", font_size=18, font_color=ACCENT_BLUE,
            bold=False, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1.2), Inches(2.9), Inches(10.5), Inches(1.8),
            "ỨNG DỤNG KIẾN TRÚC HYBRID (EDGE-CLOUD)\nVÀ DEEP LEARNING TRONG BÀI TOÁN\nNHẬN DIỆN HÀNH VI TÀI XẾ XE KHÁCH",
            font_size=28, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER,
            line_spacing=1.3)

# Info block - two columns
# Left column
info_left = [
    "Ngành: Công nghệ phần mềm",
    "Sinh viên: La Đại Lộc",
    "MSSV: 4551050116",
]
info_right = [
    "Lớp: Công nghệ thông tin 45B",
    "GVHD: TS. Nguyễn Thanh Bình",
    "Gia Lai, tháng 6 năm 2026",
]

add_multiline_textbox(slide, Inches(1.5), Inches(5.0), Inches(5), Inches(2.0),
                      info_left, font_size=16, font_color=RGBColor(0xBB, 0xDE, 0xFB),
                      alignment=PP_ALIGN.LEFT, line_spacing=1.6)
add_multiline_textbox(slide, Inches(7), Inches(5.0), Inches(5), Inches(2.0),
                      info_right, font_size=16, font_color=RGBColor(0xBB, 0xDE, 0xFB),
                      alignment=PP_ALIGN.LEFT, line_spacing=1.6)

# Bottom accent
add_shape(slide, 0, Inches(7.2), SLIDE_WIDTH, Inches(0.3), fill_color=LIGHT_BLUE)

add_speaker_notes(slide, """Em xin chào quý thầy cô trong hội đồng chấm khóa luận tốt nghiệp.
Em là La Đại Lộc, sinh viên lớp Công nghệ thông tin 45B, mã số sinh viên 4551050116.
Hôm nay em xin trình bày khóa luận tốt nghiệp với đề tài: "Ứng dụng kiến trúc Hybrid (Edge-Cloud) và Deep Learning trong bài toán nhận diện hành vi tài xế xe khách".
Đề tài được thực hiện dưới sự hướng dẫn của TS. Nguyễn Thanh Bình.
Em xin bắt đầu phần trình bày của mình.""")


# ════════════════════════════════════════════════════════════════════
# SLIDE 2: Mục lục thuyết trình
# ════════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_slide_header(slide, "MỤC LỤC THUYẾT TRÌNH")

toc_items = [
    ("01", "Bối cảnh và lý do chọn đề tài", LIGHT_BLUE),
    ("02", "Cơ sở lý thuyết và công nghệ", MEDIUM_BLUE),
    ("03", "Phân tích và thiết kế hệ thống", TEAL),
    ("04", "Xây dựng và thực nghiệm", PURPLE),
    ("05", "Kết luận và hướng phát triển", ORANGE),
]

for i, (num, title, color) in enumerate(toc_items):
    y = Inches(1.8) + Inches(i * 1.0)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.5), y, Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Title text
    add_textbox(slide, Inches(3.4), y + Inches(0.05), Inches(7), Inches(0.5),
                title, font_size=22, font_color=DARK_GRAY, bold=False)

    # Connecting line
    add_shape(slide, Inches(3.3), y + Inches(0.3), Inches(6.5), Inches(0.02),
              fill_color=RGBColor(0xE0, 0xE0, 0xE0))

add_speaker_notes(slide, """Để quý thầy cô dễ theo dõi, em xin trình bày nội dung khóa luận theo 5 phần chính.
Phần 1: Em sẽ trình bày bối cảnh và lý do chọn đề tài, bao gồm thực trạng an toàn giao thông và các hạn chế của phương pháp truyền thống.
Phần 2: Tổng quan các công nghệ sử dụng trong đề tài, bao gồm YOLO, MediaPipe Pose, SlowFast và các kỹ thuật tối ưu cho thiết bị Edge.
Phần 3: Phân tích và thiết kế kiến trúc hệ thống Hybrid Edge-Cloud.
Phần 4: Quá trình xây dựng, triển khai và đánh giá kết quả thực nghiệm.
Phần 5: Kết luận và các hướng phát triển trong tương lai.""")


# ════════════════════════════════════════════════════════════════════
# SLIDE 3: Bối cảnh và lý do chọn đề tài
# ════════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_slide_header(slide, "BỐI CẢNH VÀ LÝ DO CHỌN ĐỀ TÀI", "Chương 1: Tổng quan về bài toán và công nghệ")

# Left column - main content
add_bullet_items(slide, Inches(0.8), Inches(1.6), Inches(6.5), Inches(5.0), [
    "Tai nạn giao thông là nguyên nhân hàng đầu gây thiệt hại kinh tế và tính mạng trên toàn cầu",
    "Phần lớn tai nạn nghiêm trọng xuất phát từ yếu tố chủ quan của con người, không phải lỗi kỹ thuật",
    "Mất tập trung khi lái xe (Distracted Driving) chiếm tỷ trọng cao nhất trong chuỗi nguyên nhân",
    "Các hành vi nguy hiểm phổ biến: sử dụng điện thoại, hút thuốc, không thắt dây an toàn",
    "Cần hệ thống giám sát và cảnh báo gần thời gian thực để phòng ngừa tai nạn",
], font_size=16, spacing=1.5)

# Right side - highlight cards
add_card(slide, Inches(7.8), Inches(1.6), Inches(4.8), Inches(1.5),
         "Mất tập trung", ["Rời mắt khỏi làn đường", "Buông tay khỏi vô lăng"], "⚠", icon_bg=ORANGE)

add_card(slide, Inches(7.8), Inches(3.4), Inches(4.8), Inches(1.5),
         "Vi phạm an toàn thụ động", ["Không thắt dây an toàn", "Tăng mức độ chấn thương"], "✕", icon_bg=RED)

add_card(slide, Inches(7.8), Inches(5.2), Inches(4.8), Inches(1.5),
         "Yêu cầu cấp thiết", ["Giám sát thời gian thực", "Cảnh báo tức thời cho tài xế"], "◉", icon_bg=GREEN)

add_speaker_notes(slide, """Đề tài của em xuất phát từ thực trạng an toàn giao thông hiện nay.
Theo các báo cáo từ WHO, phần lớn các vụ tai nạn giao thông nghiêm trọng không phải do lỗi kỹ thuật của phương tiện, mà bắt nguồn từ yếu tố chủ quan của con người.
Trong đó, sự mất tập trung khi lái xe chiếm tỷ trọng cao nhất. Các hành vi nguy hiểm điển hình bao gồm: sử dụng điện thoại di động, hút thuốc và không thắt dây an toàn.
Những hành vi này làm tăng đáng kể độ trễ phản xạ, tước đi thời gian vàng để tài xế xử lý tình huống khẩn cấp.
Do đó, việc xây dựng hệ thống giám sát hành vi tài xế là rất cần thiết.""")


# ════════════════════════════════════════════════════════════════════
# SLIDE 4: Hạn chế của phương pháp truyền thống
# ════════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_slide_header(slide, "HẠN CHẾ CỦA PHƯƠNG PHÁP TRUYỀN THỐNG")

# Left column - limitations
add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.4),
            "Phương pháp hiện tại", font_size=20, font_color=RED, bold=True)

limitations = [
    "Camera hành trình chỉ ghi hình thụ động, không phân tích thời gian thực",
    "Phân tích thủ công sau khi sự cố đã xảy ra → không có khả năng cảnh báo sớm",
    "Object Detection đơn lẻ dễ gây cảnh báo sai (False Positive) do thiếu ngữ cảnh",
    "Không phân biệt được vật thể của tài xế hay hành khách ghế phụ",
    "Gây hội chứng \"mệt mỏi vì cảnh báo\" khi báo động giả liên tục",
]
add_bullet_items(slide, Inches(0.8), Inches(2.2), Inches(6), Inches(4.5),
                 limitations, font_size=15, bullet_char="✗", spacing=1.5)

# Right column - solution arrow
add_textbox(slide, Inches(7.5), Inches(1.6), Inches(5), Inches(0.4),
            "Giải pháp đề xuất", font_size=20, font_color=GREEN, bold=True)

solutions = [
    "Xử lý video thời gian thực tại Edge",
    "Kết hợp YOLO + MediaPipe Pose cho ngữ cảnh",
    "Cảnh báo tức thời tại phương tiện",
    "Kiến trúc Hybrid Edge–Cloud toàn diện",
    "Lưu vết bằng chứng kỹ thuật số",
]
add_bullet_items(slide, Inches(7.5), Inches(2.2), Inches(5), Inches(4.5),
                 solutions, font_size=15, bullet_char="✓", spacing=1.5, font_color=RGBColor(0x2E, 0x7D, 0x32))

# Arrow between columns
add_arrow(slide, Inches(6.5), Inches(3.8), Inches(0.8), Inches(0.4), fill_color=LIGHT_BLUE)

add_speaker_notes(slide, """Trước khi có công nghệ Deep Learning, các hệ thống giám sát chủ yếu dựa vào camera hành trình truyền thống.
Camera hành trình về bản chất chỉ là thiết bị ghi hình thụ động, việc phát hiện vi phạm phải thực hiện thủ công sau khi sự cố xảy ra.
Một số hệ thống tích hợp Object Detection cơ bản, nhưng chỉ phát hiện vật thể đơn lẻ mà thiếu cơ chế hiểu ngữ cảnh.
Ví dụ, hệ thống có thể nhận diện điện thoại trong khung hình, nhưng không biết đó là của tài xế hay hành khách ghế phụ.
Điều này dẫn đến tỷ lệ cảnh báo sai cao, gây mệt mỏi cho tài xế và làm giảm niềm tin vào hệ thống.""")


# ════════════════════════════════════════════════════════════════════
# SLIDE 5: Mục tiêu đề tài
# ════════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_slide_header(slide, "MỤC TIÊU ĐỀ TÀI")

# Main objective box
main_obj = add_rounded_rect(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(1.2),
                            fill_color=RGBColor(0xE3, 0xF2, 0xFD), line_color=LIGHT_BLUE)
add_textbox(slide, Inches(1.2), Inches(1.75), Inches(11), Inches(0.8),
            "Xây dựng hệ thống DMS nhận diện hành vi nguy hiểm của tài xế xe khách "
            "dựa trên kiến trúc Hybrid Edge–Cloud kết hợp Deep Learning",
            font_size=18, font_color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

# 4 objective cards
objectives = [
    ("🎯", "Phát hiện hành vi", "Nhận diện sử dụng điện thoại,\nhút thuốc, không thắt dây\nan toàn", LIGHT_BLUE),
    ("⚡", "Xử lý tại Edge", "Kết hợp YOLO + MediaPipe\nPose trên thiết bị nhúng\ntại phương tiện", MEDIUM_BLUE),
    ("☁️", "Quản lý trên Cloud", "Gửi cảnh báo và bằng chứng\nlên Cloud, hỗ trợ dashboard\nquản lý tập trung", TEAL),
    ("📊", "Đánh giá hiệu quả", "Đo lường FPS, độ chính xác\nvà khả năng truy vết\ncảnh báo", ORANGE),
]

for i, (icon, title, desc, color) in enumerate(objectives):
    x = Inches(0.8) + Inches(i * 3.05)
    y = Inches(3.2)
    card = add_rounded_rect(slide, x, y, Inches(2.8), Inches(3.2), fill_color=WHITE,
                            line_color=color)
    # Icon
    add_textbox(slide, x, y + Inches(0.2), Inches(2.8), Inches(0.5),
                icon, font_size=28, alignment=PP_ALIGN.CENTER)
    # Title
    add_textbox(slide, x + Inches(0.2), y + Inches(0.8), Inches(2.4), Inches(0.4),
                title, font_size=16, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)
    # Description
    add_textbox(slide, x + Inches(0.2), y + Inches(1.3), Inches(2.4), Inches(1.6),
                desc, font_size=13, font_color=DARK_GRAY, alignment=PP_ALIGN.CENTER,
                line_spacing=1.3)

add_speaker_notes(slide, """Đề tài có 4 mục tiêu chính:
Thứ nhất, xây dựng hệ thống DMS có khả năng nhận diện 3 nhóm hành vi nguy hiểm: sử dụng điện thoại, hút thuốc và không thắt dây an toàn.
Thứ hai, triển khai xử lý video trực tiếp tại Edge bằng cách kết hợp YOLO và MediaPipe Pose trên thiết bị nhúng.
Thứ ba, xây dựng phân hệ Cloud để tiếp nhận cảnh báo, lưu trữ bằng chứng và cung cấp dashboard giám sát tập trung.
Thứ tư, đánh giá hiệu quả hệ thống qua các chỉ số FPS, độ chính xác và khả năng truy vết cảnh báo.""")


# ════════════════════════════════════════════════════════════════════
# SLIDE 6: Phạm vi nghiên cứu
# ════════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_slide_header(slide, "PHẠM VI NGHIÊN CỨU")

# Three scope cards
scopes = [
    ("Hành vi giám sát", MEDIUM_BLUE, [
        "Sử dụng điện thoại",
        "Hút thuốc khi lái xe",
        "Không thắt dây an toàn",
    ]),
    ("Công nghệ áp dụng", TEAL, [
        "YOLO phát hiện đối tượng",
        "MediaPipe Pose ước lượng tư thế",
        "Behavior Rules Engine suy luận",
        "SlowFast xác thực hậu kiểm",
    ]),
    ("Đánh giá hiệu quả", ORANGE, [
        "FPS (Frames Per Second)",
        "Precision, Recall, F1-Score",
        "mAP50, mAP50-95",
        "Khả năng truy vết bằng chứng",
    ]),
]

for i, (title, color, items) in enumerate(scopes):
    x = Inches(0.8) + Inches(i * 4.1)
    y = Inches(1.8)

    # Header bar
    add_shape(slide, x, y, Inches(3.8), Inches(0.5), fill_color=color)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.05), Inches(3.4), Inches(0.4),
                title, font_size=18, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Content
    add_bullet_items(slide, x + Inches(0.3), y + Inches(0.7), Inches(3.2), Inches(4.0),
                     items, font_size=15, bullet_char="▸", spacing=1.5)

# Architecture highlight
add_textbox(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.8),
            "▸ Kiến trúc: Hybrid Edge–Cloud  |  ▸ Xử lý video/camera trong khoang lái  |  ▸ Triển khai nguyên mẫu thử nghiệm",
            font_size=15, font_color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

add_speaker_notes(slide, """Về phạm vi nghiên cứu, đề tài tập trung vào 3 nhóm hành vi chính: sử dụng điện thoại, hút thuốc và không thắt dây an toàn.
Về công nghệ, hệ thống sử dụng YOLO cho phát hiện đối tượng, MediaPipe Pose cho ước lượng tư thế, Behavior Rules Engine cho suy luận hành vi và SlowFast cho xác thực hậu kiểm trên Cloud.
Về đánh giá, đề tài sử dụng các chỉ số FPS để đo hiệu năng xử lý, cùng các chỉ Precision, Recall, F1-Score và mAP để đánh giá độ chính xác.
Kiến trúc hệ thống là Hybrid Edge-Cloud, xử lý video từ camera trong khoang lái.""")


# ════════════════════════════════════════════════════════════════════
# SLIDE 7: Tổng quan kiến trúc Hybrid Edge–Cloud
# ════════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_slide_header(slide, "TỔNG QUAN KIẾN TRÚC HYBRID EDGE–CLOUD", "Chương 1: Mô hình hệ thống tổng thể")

# Edge box
add_rounded_rect(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8),
                 fill_color=RGBColor(0xE3, 0xF2, 0xFD), line_color=MEDIUM_BLUE)
add_textbox(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(0.5),
            "⚙️  EDGE NODE (Thiết bị biên)", font_size=18, font_color=MEDIUM_BLUE,
            bold=True, alignment=PP_ALIGN.CENTER)

edge_items = [
    "Xử lý video tại chỗ trên NVIDIA Jetson",
    "YOLO + MediaPipe Pose suy luận thời gian thực",
    "Behavior Rules Engine đánh giá hành vi",
    "Cảnh báo cục bộ cho tài xế",
    "Gửi metadata + bằng chứng lên Cloud",
]
add_bullet_items(slide, Inches(1.2), Inches(2.6), Inches(4.8), Inches(3.5),
                 edge_items, font_size=14, bullet_char="▸", spacing=1.5)

# Arrow between
add_arrow(slide, Inches(6.5), Inches(3.5), Inches(0.8), Inches(0.5), fill_color=LIGHT_BLUE)
add_textbox(slide, Inches(6.3), Inches(4.1), Inches(1.2), Inches(0.4),
            "HTTP API", font_size=11, font_color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

# Cloud box
add_rounded_rect(slide, Inches(7.5), Inches(1.8), Inches(5.2), Inches(4.8),
                 fill_color=RGBColor(0xE8, 0xF5, 0xE9), line_color=GREEN)
add_textbox(slide, Inches(7.5), Inches(1.9), Inches(5.2), Inches(0.5),
            "☁️  CLOUD NODE (Đám mây)", font_size=18, font_color=GREEN,
            bold=True, alignment=PP_ALIGN.CENTER)

cloud_items = [
    "FastAPI tiếp nhận dữ liệu từ Edge",
    "Lưu metadata vào SQLite",
    "Lưu bằng chứng hình ảnh/video",
    "SlowFast xác thực hậu kiểm",
    "Dashboard giám sát tập trung",
]
add_bullet_items(slide, Inches(7.9), Inches(2.6), Inches(4.4), Inches(3.5),
                 cloud_items, font_size=14, bullet_char="▸", spacing=1.5)

# Bottom note
add_rounded_rect(slide, Inches(2), Inches(6.0), Inches(9), Inches(0.7),
                 fill_color=RGBColor(0xFF, 0xF3, 0xE0), line_color=ORANGE)
add_textbox(slide, Inches(2.3), Inches(6.1), Inches(8.5), Inches(0.5),
            "Nguyên tắc: \"Xử lý chiến thuật tại Biên – Quản lý chiến lược trên Đám mây\"",
            font_size=15, font_color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

add_speaker_notes(slide, """Hệ thống được thiết kế theo kiến trúc Hybrid Edge-Cloud gồm hai phân hệ chính.
Phân hệ Edge được triển khai trên thiết bị nhúng NVIDIA Jetson tại phương tiện. Edge chịu trách nhiệm tiếp nhận luồng video, chạy YOLO và MediaPipe Pose để phát hiện đối tượng và ước lượng tư thế, sau đó dùng Behavior Rules Engine để đánh giá hành vi. Khi phát hiện vi phạm, Edge kích hoạt cảnh báo cục bộ cho tài xế và gửi dữ liệu lên Cloud.
Phân hệ Cloud sử dụng FastAPI để tiếp nhận dữ liệu, lưu metadata vào SQLite, lưu bằng chứng hình ảnh/video. Cloud cũng hỗ trợ xác thực hậu kiểm bằng SlowFast và cung cấp Dashboard giám sát tập trung.
Nguyên tắc thiết kế là: Xử lý chiến thuật tại Biên, Quản lý chiến lược trên Đám mây.""")


# ════════════════════════════════════════════════════════════════════
# SLIDE 8: So sánh Edge AI và Cloud AI
# ════════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_slide_header(slide, "SO SÁNH EDGE AI VÀ CLOUD AI", "Bảng 1.1 trong báo cáo")

table_data = [
    ["Tiêu chí", "Edge AI", "Cloud AI"],
    ["Độ trễ", "< 50ms (xử lý tại chỗ)", "200-500ms+ (phụ thuộc mạng)"],
    ["Băng thông", "Thấp (chỉ gửi metadata)", "Cao (truyền video thô)"],
    ["Bảo mật", "Dữ liệu tại thiết bị", "Dữ liệu trên server (rủi ro)"],
    ["Ngoại tuyến", "Có (không cần mạng)", "Không (phụ thuộc kết nối)"],
    ["Chi phí", "Thấp", "Cao (video streaming 24/7)"],
    ["Tính toán", "Giới hạn (phần cứng nhúng)", "Lớn, có thể mở rộng"],
]

add_table(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(4.5),
          len(table_data), 3, table_data,
          col_widths=[Inches(2.5), Inches(4.5), Inches(4.5)])

# Bottom highlight
add_rounded_rect(slide, Inches(2.5), Inches(6.3), Inches(8), Inches(0.6),
                 fill_color=RGBColor(0xE3, 0xF2, 0xFD), line_color=MEDIUM_BLUE)
add_textbox(slide, Inches(2.8), Inches(6.35), Inches(7.5), Inches(0.5),
            "⟹  Kiến trúc Hybrid kết hợp ưu điểm của cả Edge AI và Cloud AI",
            font_size=16, font_color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

add_speaker_notes(slide, """Để hiểu tại sao đề tài sử dụng kiến trúc Hybrid, em xin so sánh Edge AI và Cloud AI.
Edge AI xử lý tại chỗ với độ trễ thấp dưới 50ms, chỉ cần băng thông thấp vì chỉ gửi metadata, dữ liệu vẫn nằm tại thiết bị đảm bảo bảo mật. Đặc biệt, Edge AI có thể hoạt động ngoại tuyến khi mất kết nối mạng.
Tuy nhiên, Edge AI bị giới hạn về năng lực tính toán do phần cứng nhúng.
Cloud AI có năng lực tính toán lớn và có thể mở rộng, nhưng phụ thuộc vào mạng, có độ trễ cao và chi phí vận hành lớn.
Kiến trúc Hybrid kết hợp ưu điểm của cả hai: Edge xử lý thời gian thực, Cloud quản lý và giám sát tập trung.""")


# ════════════════════════════════════════════════════════════════════
# SLIDE 9: Công nghệ sử dụng
# ════════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_slide_header(slide, "CÔNG NGHỆ SỬ DỤNG TRONG HỆ THỐNG", "Chương 2 & 3: Cơ sở lý thuyết và thiết kế")

techs = [
    ("YOLO", "Phát hiện đối tượng", "phone, smoking,\nseatbelt, no-seatbelt", MEDIUM_BLUE),
    ("MediaPipe\nPose", "Ước lượng tư thế", "13 landmarks\ncơ thể tài xế", TEAL),
    ("SlowFast", "Nhận diện hành động", "Xác thực theo\nchuỗi thời gian", PURPLE),
    ("TensorRT", "Tối ưu Edge", "Gia tốc suy luận\ntrên GPU", ORANGE),
]

for i, (name, role, desc, color) in enumerate(techs):
    x = Inches(0.8) + Inches(i * 3.05)
    y = Inches(1.8)

    # Card
    add_rounded_rect(slide, x, y, Inches(2.8), Inches(2.5), fill_color=WHITE, line_color=color)

    # Color header
    add_shape(slide, x, y, Inches(2.8), Inches(0.6), fill_color=color)
    add_textbox(slide, x, y + Inches(0.05), Inches(2.8), Inches(0.5),
                name, font_size=18, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, x + Inches(0.2), y + Inches(0.7), Inches(2.4), Inches(0.4),
                role, font_size=14, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.2), y + Inches(1.2), Inches(2.4), Inches(1.0),
                desc, font_size=13, font_color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# Cloud technologies
add_textbox(slide, Inches(0.8), Inches(4.7), Inches(12), Inches(0.4),
            "Phân hệ Cloud", font_size=18, font_color=DARK_BLUE, bold=True)

cloud_techs = [
    ("FastAPI", "Backend Framework", MEDIUM_BLUE),
    ("SQLite + SQLAlchemy", "Database + ORM", TEAL),
    ("React 19 + Vite 7", "Dashboard SPA", LIGHT_BLUE),
    ("Recharts + Lucide", "Biểu đồ + Icons", ORANGE),
]

for i, (name, desc, color) in enumerate(cloud_techs):
    x = Inches(0.8) + Inches(i * 3.05)
    y = Inches(5.3)
    add_rounded_rect(slide, x, y, Inches(2.8), Inches(1.0), fill_color=RGBColor(0xF5, 0xF5, 0xF5),
                     line_color=color)
    add_textbox(slide, x + Inches(0.1), y + Inches(0.05), Inches(2.6), Inches(0.4),
                name, font_size=15, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.1), y + Inches(0.5), Inches(2.6), Inches(0.4),
                desc, font_size=12, font_color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

add_speaker_notes(slide, """Hệ thống sử dụng các công nghệ chính sau:
YOLO cho phát hiện đối tượng, nhận diện 4 lớp: phone, smoking, seatbelt và no-seatbelt.
MediaPipe Pose để ước lượng tư thế, trích xuất 13 landmarks quan trọng của cơ thể tài xế.
SlowFast cho nhận diện hành động, hỗ trợ xác thực hành vi theo chuỗi thời gian trên Cloud.
TensorRT để tối ưu hiệu năng suy luận trên thiết bị Edge.
Tại phân hệ Cloud, hệ thống sử dụng FastAPI làm backend, SQLite kết hợp SQLAlchemy cho cơ sở dữ liệu, và React 19 với Vite 7 để xây dựng Dashboard giám sát.""")


# ════════════════════════════════════════════════════════════════════
# SLIDE 10: YOLO trong bài toán DMS
# ════════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_slide_header(slide, "YOLO TRONG BÀI TOÁN DMS", "Chương 2.1: Phát hiện đối tượng")

# Left - YOLO architecture diagram
add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.4),
            "Kiến trúc YOLO (Backbone → Neck → Head)", font_size=16, font_color=DARK_BLUE, bold=True)

# Flow diagram
boxes = [
    ("Input\n640×640", LIGHT_GRAY, DARK_GRAY),
    ("Backbone\n(CSPDarknet)", MEDIUM_BLUE, WHITE),
    ("Neck\n(FPN/PAN)", TEAL, WHITE),
    ("Head\n(Detection)", ORANGE, WHITE),
    ("Output\nBBox + Class", GREEN, WHITE),
]

for i, (text, bg, fg) in enumerate(boxes):
    x = Inches(0.8) + Inches(i * 2.2)
    add_flow_box(slide, x, Inches(2.3), Inches(1.8), Inches(1.0), text,
                 fill_color=bg, font_color=fg, font_size=12)
    if i < len(boxes) - 1:
        add_arrow(slide, x + Inches(1.85), Inches(2.6), Inches(0.3), Inches(0.3),
                  fill_color=LIGHT_BLUE)

# Detection classes
add_textbox(slide, Inches(0.8), Inches(3.6), Inches(5.5), Inches(0.4),
            "4 lớp đối tượng phát hiện:", font_size=15, font_color=DARK_BLUE, bold=True)

classes = ["phone", "smoking", "seatbelt", "no-seatbelt"]
class_colors = [LIGHT_BLUE, ORANGE, GREEN, RED]
for i, (cls, clr) in enumerate(zip(classes, class_colors)):
    x = Inches(0.8) + Inches(i * 1.5)
    add_rounded_rect(slide, x, Inches(4.1), Inches(1.3), Inches(0.5), fill_color=clr)
    add_textbox(slide, x, Inches(4.15), Inches(1.3), Inches(0.4),
                cls, font_size=13, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Right - Key points
add_textbox(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(0.4),
            "Vai trò trong hệ thống", font_size=16, font_color=DARK_BLUE, bold=True)

yolo_points = [
    "Tầng trích xuất đặc trưng không gian sơ cấp",
    "Phát hiện nhanh đối tượng liên quan đến vi phạm",
    "Đầu ra là bounding box + confidence score",
    "Kết hợp với MediaPipe Pose để xác thực ngữ cảnh",
    "Tốc độ nhanh, phù hợp thời gian thực trên Edge",
    "Mô hình: YOLO11m (Ultralytics), imgsz=768",
]
add_bullet_items(slide, Inches(7), Inches(2.2), Inches(5.5), Inches(4.5),
                 yolo_points, font_size=14, bullet_char="▸", spacing=1.4)

add_speaker_notes(slide, """YOLO là mô hình phát hiện đối tượng một giai đoạn, được đánh giá cao về tốc độ.
Kiến trúc YOLO gồm 3 phần chính: Backbone trích xuất đặc trưng, Neck kết hợp đặc trưng đa tỉ lệ, và Head đưa ra dự đoán bounding box cùng class.
Trong hệ thống, YOLO được huấn luyện để nhận diện 4 lớp đối tượng: phone, smoking, seatbelt và no-seatbelt.
YOLO đóng vai trò là tầng trích xuất đặc trưng không gian sơ cấp, cung cấp tọa độ bounding box cho các tầng xử lý phía sau.
Đầu ra của YOLO không được sử dụng trực tiếp mà kết hợp với MediaPipe Pose để xác thực ngữ cảnh, giúp giảm cảnh báo sai.
Đề tài sử dụng mô hình YOLO11m của Ultralytics với kích thước ảnh đầu vào 768x768.""")


# ════════════════════════════════════════════════════════════════════
# SLIDE 11: MediaPipe Pose
# ════════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_slide_header(slide, "MEDIAPIPE POSE – ƯỚC LƯỢNG TƯ THẾ", "Chương 2.2: 13 landmarks quan trọng")

# Left - Landmarks table
add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.4),
            "13 điểm mốc sử dụng trong hệ thống", font_size=16, font_color=DARK_BLUE, bold=True)

landmark_data = [
    ["STT", "Điểm mốc", "Vùng", "Vai trò"],
    ["1", "nose", "Mặt", "Xác định vị trí đầu"],
    ["2-3", "left/right_ear", "Mặt", "Hướng đầu"],
    ["4-5", "mouth_left/right", "Mặt", "Hút thuốc (gần miệng)"],
    ["6-7", "left/right_shoulder", "Thân trên", "Driver ROI, Chest ROI"],
    ["8-9", "left/right_elbow", "Tay", "Ngữ cảnh thao tác tay"],
    ["10-11", "left/right_wrist", "Tay", "Proximity với vật thể"],
    ["12-13", "left/right_hip", "Thân dưới", "Fallback shoulder width"],
]

add_table(slide, Inches(0.5), Inches(2.2), Inches(6.3), Inches(4.2),
          len(landmark_data), 4, landmark_data,
          col_widths=[Inches(0.7), Inches(1.8), Inches(1.2), Inches(2.6)])

# Right - Key points
add_textbox(slide, Inches(7.2), Inches(1.6), Inches(5.5), Inches(0.4),
            "Vai trò trong hệ thống", font_size=16, font_color=DARK_BLUE, bold=True)

mp_points = [
    "Cung cấp hệ tọa độ sinh trắc học của tài xế",
    "Xây dựng Driver ROI (vùng người lái)",
    "Nội suy Chest ROI (vùng ngực) cho dây an toàn",
    "Đo khoảng cách Euclidean giữa vật thể và cơ thể",
    "Hỗ trợ làm mượt theo thời gian (giảm jitter)",
    "Kiến trúc BlazePose – siêu nhẹ, phù hợp Edge",
    "Filter theo visibility ≥ 0.35 để loại nhiễu",
]
add_bullet_items(slide, Inches(7.2), Inches(2.2), Inches(5.5), Inches(4.5),
                 mp_points, font_size=14, bullet_char="▸", spacing=1.4)

add_speaker_notes(slide, """MediaPipe Pose là framework ước lượng tư thế người do Google phát triển, dựa trên kiến trúc BlazePose siêu nhẹ.
Hệ thống sử dụng 13 trong số 33 điểm mốc, tập trung vào vùng mặt, vai, tay và hông.
Các điểm mốc này được lọc theo chỉ số visibility với ngưỡng 0.35 để loại bỏ tọa độ bị nhiễu.
Vai trò chính của MediaPipe Pose trong hệ thống là: cung cấp hệ tọa độ sinh trắc học, xây dựng Driver ROI và Chest ROI, đo khoảng cách giữa vật thể và cơ thể tài xế.
Đặc biệt, MediaPipe hỗ trợ làm mượt theo thời gian, giúp giảm hiện tượng jitter của các điểm mốc giữa các khung hình liên tiếp.""")


# ════════════════════════════════════════════════════════════════════
# SLIDE 12: SlowFast
# ════════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_slide_header(slide, "MÔ HÌNH SLOWFAST – NHẬN DIỆN HÀNH ĐỘNG", "Chương 2.3: Phân tích không gian – thời gian")

# Two pathway diagram
add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.4),
            "Kiến trúc 2-Pathway: Slow Pathway + Fast Pathway", font_size=16, font_color=DARK_BLUE, bold=True)

# Slow pathway
add_rounded_rect(slide, Inches(0.8), Inches(2.3), Inches(5.5), Inches(2.0),
                 fill_color=RGBColor(0xE3, 0xF2, 0xFD), line_color=MEDIUM_BLUE)
add_textbox(slide, Inches(0.8), Inches(2.4), Inches(5.5), Inches(0.4),
            "🐢 Slow Pathway", font_size=18, font_color=MEDIUM_BLUE, bold=True,
            alignment=PP_ALIGN.CENTER)
slow_items = [
    "Lấy mẫu thưa hơn (ít frame hơn)",
    "Học đặc trưng ngữ nghĩa không gian ổn định",
    "Bố cục cảnh, hình dạng vật thể, cấu trúc cơ thể",
]
add_bullet_items(slide, Inches(1.2), Inches(3.0), Inches(4.8), Inches(1.2),
                 slow_items, font_size=13, bullet_char="▸", spacing=1.3)

# Fast pathway
add_rounded_rect(slide, Inches(6.8), Inches(2.3), Inches(5.5), Inches(2.0),
                 fill_color=RGBColor(0xFF, 0xF3, 0xE0), line_color=ORANGE)
add_textbox(slide, Inches(6.8), Inches(2.4), Inches(5.5), Inches(0.4),
            "🐇 Fast Pathway", font_size=18, font_color=ORANGE, bold=True,
            alignment=PP_ALIGN.CENTER)
fast_items = [
    "Lấy mẫu dày hơn (nhiều frame hơn)",
    "Tập trung thông tin thời gian, chuyển động nhanh",
    "Thiết kế nhẹ hơn về kênh đặc trưng",
]
add_bullet_items(slide, Inches(7.2), Inches(3.0), Inches(4.8), Inches(1.2),
                 fast_items, font_size=13, bullet_char="▸", spacing=1.3)

# Lateral connections
add_shape(slide, Inches(5.5), Inches(2.8), Inches(2.0), Inches(0.8),
          fill_color=RGBColor(0xF3, 0xE5, 0xF5), line_color=PURPLE)
add_textbox(slide, Inches(5.5), Inches(2.9), Inches(2.0), Inches(0.6),
            "Lateral\nConnections", font_size=12, font_color=PURPLE, bold=True,
            alignment=PP_ALIGN.CENTER)

# Bottom - application in DMS
add_textbox(slide, Inches(0.8), Inches(4.8), Inches(11.5), Inches(0.4),
            "Ứng dụng trong kiến trúc hệ thống", font_size=16, font_color=DARK_BLUE, bold=True)

app_items = [
    "Triển khai tại Cloud (không chạy trên Edge do chi phí tính toán cao)",
    "Vai trò: tầng xác thực chuyên sâu (Deep Verification Layer)",
    "Phân tích video clip bằng chứng được Edge gửi lên",
    "Hiện tại sử dụng mô hình pretrained trên Kinetics-400, chưa fine-tune chuyên sâu",
    "Kích hoạt theo yêu cầu qua API POST /alerts/{id}/verify",
]
add_bullet_items(slide, Inches(0.8), Inches(5.3), Inches(11.5), Inches(1.5),
                 app_items, font_size=14, bullet_char="▸", spacing=1.3)

add_speaker_notes(slide, """SlowFast là kiến trúc mạng nơ-ron cho phân tích video, gồm hai pathway chạy song song.
Slow Pathway xử lý ít frame hơn, tập trung học đặc trưng không gian như hình dạng, cấu trúc cơ thể.
Fast Pathway xử lý nhiều frame hơn, tập trung vào thông tin thời gian và chuyển động nhanh.
Hai pathway được kết nối qua lateral connections để trao đổi thông tin.
Trong kiến trúc hệ thống, SlowFast được triển khai tại Cloud, đóng vai trò tầng xác thực chuyên sâu.
Khi Edge phát hiện nghi vấn, video clip được gửi lên Cloud để SlowFast phân tích và xác nhận.
Hiện tại, hệ thống sử dụng mô hình SlowFast đã huấn luyện sẵn trên Kinetics-400, chưa fine-tune chuyên sâu cho hành vi tài xế.""")


# ════════════════════════════════════════════════════════════════════
# SLIDE 13: Tối ưu mô hình cho Edge
# ════════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_slide_header(slide, "TỐI ƯU MÔ HÌNH CHO THIẾT BỊ EDGE", "Chương 2.4: Quantization, TensorRT, Frame Skipping")

# Three optimization cards
opt_data = [
    ("Lượng tử hóa (Quantization)", MEDIUM_BLUE, [
        "Chuyển FP32 → FP16/INT8",
        "Giảm 2-4 lần kích thước mô hình",
        "Tăng thông lượng tính toán",
        "Post-Training Quantization (PTQ)",
        "Quantization-Aware Training (QAT)",
    ]),
    ("TensorRT", TEAL, [
        "Trình biên dịch của NVIDIA",
        "Hợp nhất lớp đồ thị (Layer Fusion)",
        "Tự động tinh chỉnh nhân tính toán",
        "Tối ưu bộ nhớ động",
        "Hỗ trợ FP16/INT8 với Tensor Cores",
    ]),
    ("Chiến lược runtime", ORANGE, [
        "Frame Skipping (bỏ khung có chủ đích)",
        "detect_every_n_frames = 2",
        "pose_every_n_frames = 3",
        "Resize đầu vào 640×640",
        "Kích hoạt Pose theo điều kiện (lazy)",
    ]),
]

for i, (title, color, items) in enumerate(opt_data):
    x = Inches(0.8) + Inches(i * 4.1)
    y = Inches(1.8)

    add_shape(slide, x, y, Inches(3.8), Inches(0.5), fill_color=color)
    add_textbox(slide, x + Inches(0.1), y + Inches(0.05), Inches(3.6), Inches(0.4),
                title, font_size=15, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_bullet_items(slide, x + Inches(0.3), y + Inches(0.7), Inches(3.2), Inches(4.0),
                     items, font_size=13, bullet_char="▸", spacing=1.4)

# Bottom note
add_rounded_rect(slide, Inches(1.5), Inches(6.0), Inches(10), Inches(0.7),
                 fill_color=RGBColor(0xFF, 0xF3, 0xE0), line_color=ORANGE)
add_textbox(slide, Inches(1.8), Inches(6.1), Inches(9.5), Inches(0.5),
            "Mục tiêu: Đảm bảo FPS phục vụ cảnh báo gần thời gian thực trên thiết bị Edge có tài nguyên hạn chế",
            font_size=14, font_color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

add_speaker_notes(slide, """Để triển khai trên thiết bị Edge có tài nguyên hạn chế, đề tài áp dụng 3 nhóm kỹ thuật tối ưu.
Thứ nhất là lượng tử hóa, chuyển đổi từ FP32 xuống FP16 hoặc INT8, giúp giảm kích thước mô hình và tăng tốc tính toán.
Thứ hai là TensorRT, trình biên dịch của NVIDIA giúp hợp nhất các lớp tính toán, tự động tinh chỉnh nhân và tối ưu bộ nhớ.
Thứ ba là các chiến lược runtime như Frame Skipping, chỉ chạy YOLO mỗi 2 frame và MediaPipe mỗi 3 frame, cùng với việc resize đầu vào về 640 pixel.
Đặc biệt, MediaPipe Pose chỉ được kích hoạt khi YOLO phát hiện đối tượng khả nghi, giúp tiết kiệm tài nguyên đáng kể.""")


# ════════════════════════════════════════════════════════════════════
# SLIDE 14: Pipeline xử lý tại Edge
# ════════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_slide_header(slide, "PIPELINE XỬ LÝ TẠI EDGE", "Chương 3.2.1: Luồng xử lý tuần tự")

# Horizontal flow pipeline
pipeline_steps = [
    ("Camera\nInput", DARK_GRAY),
    ("Resize\n640px", GRAY),
    ("Frame\nCounter", MEDIUM_BLUE),
    ("YOLO\nDetection", LIGHT_BLUE),
    ("Check\nObjects", TEAL),
    ("MediaPipe\nPose", GREEN),
    ("Rules\nEngine", ORANGE),
    ("Violation\nEvent", RED),
]

for i, (text, color) in enumerate(pipeline_steps):
    x = Inches(0.4) + Inches(i * 1.58)
    y = Inches(1.8)
    add_flow_box(slide, x, y, Inches(1.35), Inches(0.9), text,
                 fill_color=color, font_color=WHITE, font_size=11)
    if i < len(pipeline_steps) - 1:
        add_arrow(slide, x + Inches(1.38), y + Inches(0.25), Inches(0.18), Inches(0.3),
                  fill_color=LIGHT_BLUE)

# Output branches
add_down_arrow(slide, Inches(10.5), Inches(2.8), Inches(0.4), Inches(0.5), fill_color=RED)

outputs = [
    ("🔊 AlertManager", "Cảnh báo cục bộ\n(âm thanh, LED)", Inches(8.0)),
    ("💾 EvidenceWriter", "Lưu ảnh/video\nbằng chứng", Inches(10.0)),
    ("🖥️ OverlayRenderer", "Hiển thị kết quả\ntrên màn hình", Inches(12.0)),
]

for title, desc, x in outputs:
    add_rounded_rect(slide, x, Inches(3.5), Inches(2.2), Inches(1.2),
                     fill_color=RGBColor(0xFC, 0xE4, 0xEC), line_color=RED)
    add_textbox(slide, x + Inches(0.1), Inches(3.55), Inches(2.0), Inches(0.35),
                title, font_size=12, font_color=RED, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.1), Inches(3.95), Inches(2.0), Inches(0.6),
                desc, font_size=11, font_color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# Key optimization note
add_textbox(slide, Inches(0.8), Inches(5.2), Inches(12), Inches(0.4),
            "Ba cơ chế tối ưu chính tại Edge:", font_size=16, font_color=DARK_BLUE, bold=True)

opt_notes = [
    ("Hạ độ phân giải", "Resize từ 848→640px", Inches(0.8)),
    ("Frame Skipping", "YOLO mỗi N=2 frame", Inches(4.5)),
    ("Lazy Pose", "Chỉ chạy Pose khi có detection", Inches(8.2)),
]

for title, desc, x in opt_notes:
    add_rounded_rect(slide, x, Inches(5.7), Inches(3.3), Inches(0.9),
                     fill_color=RGBColor(0xE3, 0xF2, 0xFD), line_color=MEDIUM_BLUE)
    add_textbox(slide, x + Inches(0.1), Inches(5.75), Inches(3.1), Inches(0.35),
                title, font_size=14, font_color=MEDIUM_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.1), Inches(6.1), Inches(3.1), Inches(0.35),
                desc, font_size=12, font_color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

add_speaker_notes(slide, """Pipeline xử lý tại Edge được thiết kế tuần tự gồm 8 bước chính.
Bắt đầu từ Camera Input, khung hình được resize về 640 pixel, sau đó Frame Counter kiểm soát tần suất chạy mô hình.
YOLO Detection chỉ chạy mỗi 2 frame để tiết kiệm tài nguyên. Nếu phát hiện đối tượng liên quan, MediaPipe Pose được kích hoạt.
Kết quả từ cả hai mô hình được đưa vào Behavior Rules Engine để đánh giá hành vi vi phạm.
Khi xác nhận vi phạm, hệ thống đồng thời: kích hoạt cảnh báo cục bộ qua AlertManager, lưu bằng chứng qua EvidenceWriter, và hiển thị kết quả qua OverlayRenderer.
Ba cơ chế tối ưu chính là: hạ độ phân giải, frame skipping và lazy activation cho Pose.""")


# ════════════════════════════════════════════════════════════════════
# SLIDE 15: Cơ chế kết hợp YOLO và MediaPipe Pose
# ════════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_slide_header(slide, "CƠ CHẾ KẾT HỢP YOLO VÀ MEDIAPIPE POSE", "Chương 3.2.2: Suy luận hai tầng")

# Left - Flow diagram
add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.4),
            "Luồng suy luận hai tầng", font_size=16, font_color=DARK_BLUE, bold=True)

# YOLO layer
add_rounded_rect(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(1.2),
                 fill_color=RGBColor(0xE3, 0xF2, 0xFD), line_color=MEDIUM_BLUE)
add_textbox(slide, Inches(0.8), Inches(2.25), Inches(5.5), Inches(0.35),
            "Tầng 1: YOLO Detection", font_size=14, font_color=MEDIUM_BLUE, bold=True,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.0), Inches(2.65), Inches(5.0), Inches(0.6),
            "Phát hiện đối tượng → Bounding box + Confidence\nCung cấp tín hiệu kích hoạt cho Pose",
            font_size=12, font_color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

add_down_arrow(slide, Inches(3.3), Inches(3.4), Inches(0.4), Inches(0.3), fill_color=LIGHT_BLUE)

# Pose layer
add_rounded_rect(slide, Inches(0.8), Inches(3.8), Inches(5.5), Inches(1.2),
                 fill_color=RGBColor(0xE8, 0xF5, 0xE9), line_color=GREEN)
add_textbox(slide, Inches(0.8), Inches(3.85), Inches(5.5), Inches(0.35),
            "Tầng 2: MediaPipe Pose", font_size=14, font_color=GREEN, bold=True,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.0), Inches(4.25), Inches(5.0), Inches(0.6),
            "13 landmarks → Driver ROI + Chest ROI\nĐo khoảng cách Euclidean với vật thể",
            font_size=12, font_color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

add_down_arrow(slide, Inches(3.3), Inches(5.0), Inches(0.4), Inches(0.3), fill_color=GREEN)

# Rules Engine
add_rounded_rect(slide, Inches(0.8), Inches(5.4), Inches(5.5), Inches(1.0),
                 fill_color=RGBColor(0xFF, 0xF3, 0xE0), line_color=ORANGE)
add_textbox(slide, Inches(0.8), Inches(5.45), Inches(5.5), Inches(0.35),
            "Tầng 3: Behavior Rules Engine", font_size=14, font_color=ORANGE, bold=True,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.0), Inches(5.85), Inches(5.0), Inches(0.4),
            "Score = w1×YOLO_conf + w2×proximity(obj,wrist) + w3×proximity(obj,face) − w4×size_penalty",
            font_size=11, font_color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# Right - Key principles
add_textbox(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(0.4),
            "Nguyên tắc xác thực", font_size=16, font_color=DARK_BLUE, bold=True)

principles = [
    "Không fallback YOLO đơn lẻ cho hành vi dùng điện thoại",
    "Fallback cho hút thuốc: raw conf ≥ 0.70 khi không có pose",
    "Cần tối thiểu: 1 điểm mặt + 1 điểm vai + 1 điểm cổ tay",
    "Khoảng cách chuẩn hóa theo tỷ lệ vai tài xế",
    "Ngưỡng xác nhận: score ≥ 0.62",
    "Số frame liên tiếp cần xác nhận: 7 frame (phone, smoking), 12 frame (no_seatbelt)",
]
add_bullet_items(slide, Inches(7), Inches(2.2), Inches(5.8), Inches(4.5),
                 principles, font_size=13, bullet_char="▸", spacing=1.4)

# Benefit highlight
add_rounded_rect(slide, Inches(7), Inches(5.8), Inches(5.5), Inches(0.8),
                 fill_color=RGBColor(0xE8, 0xF5, 0xE9), line_color=GREEN)
add_textbox(slide, Inches(7.3), Inches(5.9), Inches(5.0), Inches(0.6),
            "✓ Giảm cảnh báo sai tốt hơn so với chỉ dùng YOLO thuần túy",
            font_size=14, font_color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

add_speaker_notes(slide, """Hệ thống sử dụng cơ chế suy luận hai tầng kết hợp YOLO và MediaPipe Pose.
Tầng 1 là YOLO, phát hiện đối tượng và cung cấp bounding box cùng confidence score.
Tầng 2 là MediaPipe Pose, trích xuất 13 landmarks để xây dựng Driver ROI và Chest ROI, đo khoảng cách giữa vật thể và cơ thể.
Tầng 3 là Behavior Rules Engine, tính điểm vi phạm dựa trên công thức kết hợp confidence của YOLO, khoảng cách đến cổ tay và vùng mặt.
Đối với hành vi dùng điện thoại, hệ thống không fallback về YOLO đơn lẻ nếu không có pose.
Đối với hút thuốc, có cơ chế fallback với confidence thô ≥ 0.70 khi pose không khả dụng.
Cơ chế này giúp giảm đáng kể cảnh báo sai so với phương pháp chỉ dùng YOLO.""")


# ════════════════════════════════════════════════════════════════════
# SLIDE 16: Driver ROI và Chest ROI
# ════════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_slide_header(slide, "DRIVER ROI VÀ CHEST ROI", "Chương 3.2.2: Vùng quan tâm trong khoang lái")

# Left - Driver ROI
add_rounded_rect(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8),
                 fill_color=RGBColor(0xE3, 0xF2, 0xFD), line_color=MEDIUM_BLUE)
add_textbox(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(0.4),
            "Driver ROI (Vùng người lái)", font_size=18, font_color=MEDIUM_BLUE,
            bold=True, alignment=PP_ALIGN.CENTER)

driver_roi_items = [
    "Xây dựng từ tọa độ: mũi, tai, vai, cổ tay",
    "Giới hạn vùng thao tác chính của tài xế",
    "Phân biệt vật thể của tài xế vs hành khách",
    "Giảm nhiễu từ ghế phụ và bảng điều khiển",
    "Tập trung tính toán Edge vào đúng đối tượng",
]
add_bullet_items(slide, Inches(1.2), Inches(2.5), Inches(4.8), Inches(3.5),
                 driver_roi_items, font_size=14, bullet_char="▸", spacing=1.5)

# Right - Chest ROI
add_rounded_rect(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(4.8),
                 fill_color=RGBColor(0xFF, 0xF3, 0xE0), line_color=ORANGE)
add_textbox(slide, Inches(7), Inches(1.9), Inches(5.5), Inches(0.4),
            "Chest ROI (Vùng ngực)", font_size=18, font_color=ORANGE,
            bold=True, alignment=PP_ALIGN.CENTER)

chest_roi_items = [
    "Nội suy từ tọa độ hai vai và trục xương sống",
    "Phục vụ đánh giá trạng thái dây an toàn",
    "Đoán quỹ đạo chuẩn mà dây phải đi qua",
    "Đối chiếu bounding box dây an toàn với ranh giới",
    "Loại bỏ báo động giả do nếp áo hoặc quai túi",
]
add_bullet_items(slide, Inches(7.4), Inches(2.5), Inches(4.8), Inches(3.5),
                 chest_roi_items, font_size=14, bullet_char="▸", spacing=1.5)

# Bottom summary
add_rounded_rect(slide, Inches(2), Inches(6.0), Inches(9), Inches(0.7),
                 fill_color=RGBColor(0xF3, 0xE5, 0xF5), line_color=PURPLE)
add_textbox(slide, Inches(2.3), Inches(6.1), Inches(8.5), Inches(0.5),
            "Cả hai ROI đều được xây dựng từ dữ liệu MediaPipe Pose → giảm phụ thuộc vào góc camera",
            font_size=14, font_color=PURPLE, bold=True, alignment=PP_ALIGN.CENTER)

add_speaker_notes(slide, """Hệ thống xây dựng hai vùng quan tâm quan trọng từ dữ liệu MediaPipe Pose.
Driver ROI là vùng thao tác chính của tài xế, được xây dựng từ tọa độ mũi, tai, vai và cổ tay.
Vùng này giúp hệ thống phân biệt vật thể của tài xế với vật thể của hành khách hoặc trên bảng điều khiển.
Chest ROI là vùng ngực, được nội suy từ hai vai và trục xương sống.
Vùng này phục vụ đánh giá trạng thái dây an toàn, bằng cách đối chiếu bounding box dây an toàn với ranh giới Chest ROI.
Cả hai ROI đều được xây dựng động từ dữ liệu Pose, giúp giảm phụ thuộc vào góc camera và vị trí lắp đặt.""")


# ════════════════════════════════════════════════════════════════════
# SLIDE 17: Phát hiện hành vi sử dụng điện thoại
# ════════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_slide_header(slide, "PHÁT HIỆN HÀNH VI SỬ DỤNG ĐIỆN THOẠI", "Chương 3.2.3")

# Flow diagram
flow_steps = [
    ("YOLO phát hiện\nđiện thoại", LIGHT_BLUE),
    ("Kiểm tra trong\nDriver ROI", MEDIUM_BLUE),
    ("Tính khoảng cách\nvới cổ tay/mặt", TEAL),
    ("Tính điểm\nvi phạm", ORANGE),
    ("Score ≥ 0.62\n× 7 frames?", RED),
    ("Cảnh báo\nvi phạm", GREEN),
]

for i, (text, color) in enumerate(flow_steps):
    x = Inches(0.5) + Inches(i * 2.1)
    add_flow_box(slide, x, Inches(1.8), Inches(1.85), Inches(1.0), text,
                 fill_color=color, font_color=WHITE, font_size=11)
    if i < len(flow_steps) - 1:
        add_arrow(slide, x + Inches(1.88), Inches(2.1), Inches(0.2), Inches(0.3),
                  fill_color=LIGHT_BLUE)

# Detail points
add_textbox(slide, Inches(0.8), Inches(3.3), Inches(11.5), Inches(0.4),
            "Chi tiết cơ chế phát hiện:", font_size=16, font_color=DARK_BLUE, bold=True)

phone_details = [
    "YOLO khoanh vùng điện thoại trong khung hình với confidence score",
    "Kiểm tra bounding box có nằm trong Driver ROI hay không → loại bỏ nếu ngoài vùng",
    "Tính khoảng cách Euclidean từ tâm điện thoại đến cổ tay và vùng đầu (tai/mũi)",
    "Điểm vi phạm = w1×YOLO_conf + w2×proximity(obj,wrist) + w3×proximity(obj,face) − w4×size_penalty",
    "Khoảng cách được chuẩn hóa theo tỷ lệ vai tài xế để phù hợp với nhiều vóc dáng",
    "Yêu cầu xác nhận liên tục 7 frame vượt ngưỡng 0.62 trước khi tạo cảnh báo",
    "Áp dụng cơ chế cooldown 4 giây giữa 2 cảnh báo cùng loại để tránh spam",
]
add_bullet_items(slide, Inches(0.8), Inches(3.8), Inches(11.5), Inches(3.0),
                 phone_details, font_size=14, bullet_char="▸", spacing=1.3)

add_speaker_notes(slide, """Đối với hành vi sử dụng điện thoại, hệ thống thực hiện theo luồng 6 bước.
Đầu tiên, YOLO phát hiện điện thoại trong khung hình. Sau đó kiểm tra xem bounding box có nằm trong Driver ROI không.
Nếu nằm trong vùng người lái, hệ thống tính khoảng cách Euclidean từ điện thoại đến cổ tay và vùng đầu.
Điểm vi phạm được tính theo công thức kết hợp confidence của YOLO, proximity với cổ tay và vùng mặt.
Hệ thống yêu cầu xác nhận liên tục 7 frame vượt ngưỡng 0.62 trước khi tạo cảnh báo, giúp tránh cảnh báo sai do nhiễu tạm thời.
Ngoài ra, cơ chế cooldown 4 giây giữa 2 cảnh báo cùng loại giúp tránh spam cảnh báo.""")


# ════════════════════════════════════════════════════════════════════
# SLIDE 18: Phát hiện hút thuốc
# ════════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_slide_header(slide, "PHÁT HIỆN HÀNH VI HÚT THUỐC", "Chương 3.2.3")

# Main content - two columns
add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.4),
            "Đặc thù bài toán", font_size=16, font_color=RED, bold=True)

smoking_challenges = [
    "Điếu thuốc có kích thước pixel rất nhỏ",
    "Dễ bị che khuất bởi ngón tay hoặc khuôn mặt",
    "Vật thể tương đồng gây nhiễu: bút, ống hút, tăm",
    "Cần kết hợp vị trí miệng, tay và tư thế",
]
add_bullet_items(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(2.5),
                 smoking_challenges, font_size=14, bullet_char="⚠", spacing=1.4)

add_textbox(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(0.4),
            "Cơ chế phát hiện", font_size=16, font_color=GREEN, bold=True)

smoking_mechanism = [
    "YOLO phát hiện vùng nghi vấn chứa điếu thuốc",
    "MediaPipe kiểm tra khoảng cách cổ tay ↔ miệng",
    "Score ≥ 0.62 trong 7 frame liên tiếp",
    "Có cơ chế fallback: raw conf ≥ 0.70 khi không có pose",
]
add_bullet_items(slide, Inches(7), Inches(2.2), Inches(5.5), Inches(2.5),
                 smoking_mechanism, font_size=14, bullet_char="✓", spacing=1.4)

# Fallback mechanism highlight
add_rounded_rect(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(1.5),
                 fill_color=RGBColor(0xFF, 0xF3, 0xE0), line_color=ORANGE)
add_textbox(slide, Inches(1.2), Inches(5.1), Inches(11), Inches(0.4),
            "Cơ chế Fallback cho hút thuốc:", font_size=16, font_color=ORANGE, bold=True)
add_textbox(slide, Inches(1.2), Inches(5.5), Inches(11), Inches(0.8),
            "Khi MediaPipe Pose không khả dụng (thất bại phát hiện landmarks), hệ thống vẫn có thể xác nhận "
            "hành vi hút thuốc dựa trên confidence thô của YOLO nếu giá trị ≥ 0.70. Cơ chế này giúp tránh "
            "bỏ sót trong điều kiện pose thất bại, tuy nhiên chỉ áp dụng cho hành vi hút thuốc.",
            font_size=13, font_color=DARK_GRAY, line_spacing=1.4)

add_speaker_notes(slide, """Hành vi hút thuốc là thách thức lớn nhất trong 3 hành vi cần phát hiện.
Điếu thuốc có kích thước rất nhỏ, dễ bị che khuất bởi ngón tay, và có nhiều vật thể tương đồng gây nhiễu như bút, ống hút.
Hệ thống phát hiện hút thuốc bằng cách: YOLO xác định vùng nghi vấn, sau đó MediaPipe kiểm tra khoảng cách giữa cổ tay và vùng miệng.
Đặc biệt, hành vi hút thuốc có cơ chế fallback: khi MediaPipe Pose không khả dụng, hệ thống vẫn có thể xác nhận dựa trên confidence thô của YOLO nếu ≥ 0.70.
Cơ chế fallback này giúp tránh bỏ sót trong điều kiện pose thất bại, nhưng chỉ áp dụng riêng cho hành vi hút thuốc.""")


# ════════════════════════════════════════════════════════════════════
# SLIDE 19: Phát hiện không thắt dây an toàn
# ════════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_slide_header(slide, "PHÁT HIỆN KHÔNG THẮT DÂY AN TOÀN", "Chương 3.2.3")

# Left content
add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.4),
            "Cơ chế phát hiện", font_size=16, font_color=DARK_BLUE, bold=True)

seatbelt_items = [
    "Sử dụng vùng Chest ROI (nội suy từ hai vai)",
    "YOLO phát hiện trạng thái seatbelt / no-seatbelt",
    "Đối chiếu bounding box với ranh giới Chest ROI",
    "Kiểm tra sự tồn tại của dải dây đi cắt ngang qua vùng ngực",
    "Ngưỡng tin cậy: 0.45 + margin 0.07",
    "Yêu cầu xác nhận 12 frame liên tiếp",
    "Không có cơ chế fallback (cần Pose để xác định Chest ROI)",
]
add_bullet_items(slide, Inches(0.8), Inches(2.2), Inches(6), Inches(4.5),
                 seatbelt_items, font_size=14, bullet_char="▸", spacing=1.4)

# Right - Importance card
add_rounded_rect(slide, Inches(7.5), Inches(1.8), Inches(5), Inches(4.5),
                 fill_color=RGBColor(0xFC, 0xE4, 0xEC), line_color=RED)
add_textbox(slide, Inches(7.5), Inches(1.9), Inches(5), Inches(0.4),
            "🛡️ Ý nghĩa an toàn thụ động", font_size=16, font_color=RED,
            bold=True, alignment=PP_ALIGN.CENTER)

importance_items = [
    "Không trực tiếp gây ra tai nạn",
    "Nhưng tước đi lớp bảo vệ cơ học quan trọng",
    "Khuếch đại mức độ nghiêm trọng khi sự cố xảy ra",
    "Cảnh báo giúp tài xế tuân thủ quy tắc an toàn",
    "Đặc biệt quan trọng trên xe khách đường dài",
]
add_bullet_items(slide, Inches(7.9), Inches(2.5), Inches(4.3), Inches(3.5),
                 importance_items, font_size=14, bullet_char="•", spacing=1.5, font_color=RED)

# Bottom note
add_rounded_rect(slide, Inches(2), Inches(6.0), Inches(9), Inches(0.7),
                 fill_color=RGBColor(0xE8, 0xF5, 0xE9), line_color=GREEN)
add_textbox(slide, Inches(2.3), Inches(6.1), Inches(8.5), Inches(0.5),
            "Thách thức: dây có thể hòa lẫn với áo, quai túi; cần dataset đa dạng về trang phục",
            font_size=14, font_color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

add_speaker_notes(slide, """Đối với hành vi không thắt dây an toàn, hệ thống sử dụng vùng Chest ROI được xây dựng từ MediaPipe Pose.
YOLO phát hiện trạng thái seatbelt hoặc no-seatbelt, sau đó đối chiếu bounding box với ranh giới Chest ROI.
Hệ thống kiểm tra sự tồn tại của dải dây đi cắt ngang qua vùng ngực. Nếu không phát hiện được dải dây, cảnh báo sẽ được tạo.
Ngưỡng tin cậy cho phát hiện dây an toàn là 0.45 với margin 0.07, yêu cầu xác nhận 12 frame liên tiếp.
Đây là hành vi an toàn thụ động, không trực tiếp gây tai nạn nhưng tước đi lớp bảo vệ quan trọng khi sự cố xảy ra.
Thách thức lớn nhất là dây an toàn có thể hòa lẫn với áo hoặc quai túi, cần dataset đa dạng về trang phục.""")


# ════════════════════════════════════════════════════════════════════
# SLIDE 20: Thiết kế Cloud Backend
# ════════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_slide_header(slide, "THIẾT KẾ CLOUD BACKEND", "Chương 3.3: Kiến trúc phân lớp")

# Three layer diagram
layers = [
    ("Router Layer", "API Routes", "Định nghĩa endpoints\nTiếp nhận request từ Edge\nvà Dashboard", MEDIUM_BLUE),
    ("Schema Layer", "Pydantic Validation", "Kiểm tra dữ liệu đầu vào\nChuẩn hóa metadata\nĐảm bảo cấu trúc thống nhất", TEAL),
    ("Model Layer", "SQLAlchemy ORM", "Ánh xạ bảng DB thành đối tượng\nThao tác CRUD\nQuản lý bảng Alerts", ORANGE),
]

for i, (title, subtitle, desc, color) in enumerate(layers):
    x = Inches(0.8) + Inches(i * 4.1)
    y = Inches(1.8)

    # Layer box
    add_rounded_rect(slide, x, y, Inches(3.8), Inches(3.5), fill_color=WHITE, line_color=color)

    # Header
    add_shape(slide, x, y, Inches(3.8), Inches(0.8), fill_color=color)
    add_textbox(slide, x, y + Inches(0.05), Inches(3.8), Inches(0.35),
                title, font_size=18, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, y + Inches(0.4), Inches(3.8), Inches(0.35),
                subtitle, font_size=12, font_color=RGBColor(0xE0, 0xE0, 0xE0),
                alignment=PP_ALIGN.CENTER)

    # Description
    add_textbox(slide, x + Inches(0.3), y + Inches(1.0), Inches(3.2), Inches(2.2),
                desc, font_size=14, font_color=DARK_GRAY, alignment=PP_ALIGN.CENTER,
                line_spacing=1.4)

# API endpoints
add_textbox(slide, Inches(0.8), Inches(5.6), Inches(12), Inches(0.4),
            "API Endpoints chính:", font_size=16, font_color=DARK_BLUE, bold=True)

apis = [
    ("POST /alerts", "Edge gửi cảnh báo (multipart/form-data)", MEDIUM_BLUE),
    ("GET /api/alerts", "Dashboard truy vấn danh sách cảnh báo", TEAL),
    ("POST /alerts/{id}/verify", "Kích hoạt xác thực SlowFast", ORANGE),
]

for i, (endpoint, desc, color) in enumerate(apis):
    x = Inches(0.8) + Inches(i * 4.1)
    y = Inches(6.1)
    add_rounded_rect(slide, x, y, Inches(3.8), Inches(0.6), fill_color=RGBColor(0xF5, 0xF5, 0xF5),
                     line_color=color)
    add_textbox(slide, x + Inches(0.1), y + Inches(0.02), Inches(3.6), Inches(0.25),
                endpoint, font_size=13, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.1), y + Inches(0.3), Inches(3.6), Inches(0.25),
                desc, font_size=11, font_color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

add_speaker_notes(slide, """Backend Cloud được thiết kế theo kiến trúc phân lớp gồm 3 lớp chính.
Router Layer định nghĩa các API endpoints, tiếp nhận request từ Edge và Dashboard.
Schema Layer sử dụng Pydantic để kiểm tra và chuẩn hóa dữ liệu đầu vào, đảm bảo cấu trúc thống nhất.
Model Layer sử dụng SQLAlchemy ORM để ánh xạ cơ sở dữ liệu, thực hiện các thao tác CRUD.
Hệ thống có 3 API endpoints chính: POST /alerts để Edge gửi cảnh báo, GET /api/alerts để Dashboard truy vấn, và POST /alerts/{id}/verify để kích hoạt xác thực SlowFast.""")


# ════════════════════════════════════════════════════════════════════
# SLIDE 21: Cơ sở dữ liệu và API
# ════════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_slide_header(slide, "CƠ SỞ DỮ LIỆU VÀ API", "Chương 3.4.1: Bảng Alerts")

# Database table
db_data = [
    ["Trường", "Mô tả", "Loại"],
    ["id", "Khóa chính (auto-increment)", "Integer PK"],
    ["event_type", "Loại vi phạm", "String"],
    ["timestamp", "Thời gian (ISO string)", "String"],
    ["confidence", "Độ tin cậy", "Float"],
    ["frame_index", "Chỉ số frame", "Integer"],
    ["source_device", "Mã thiết bị", "String"],
    ["frame_path", "Đường dẫn ảnh bằng chứng", "String"],
    ["clip_path", "Đường dẫn video clip", "String"],
    ["verified", "Trạng thái xác thực", "Boolean"],
    ["review_status", "Trạng thái kiểm duyệt", "String"],
    ["created_at", "Thời điểm tạo", "DateTime"],
]

add_table(slide, Inches(0.5), Inches(1.6), Inches(7.5), Inches(5.0),
          len(db_data), 3, db_data,
          col_widths=[Inches(2.0), Inches(3.5), Inches(2.0)])

# Right - Data flow
add_textbox(slide, Inches(8.5), Inches(1.6), Inches(4.3), Inches(0.4),
            "Luồng dữ liệu tại Cloud", font_size=16, font_color=DARK_BLUE, bold=True)

flow_items = [
    "Edge gửi multipart/form-data",
    "POST /alerts endpoint tiếp nhận",
    "Kiểm tra tính hợp lệ (Pydantic)",
    "Lưu file vào File System",
    "Ghi metadata vào SQLite",
    "Trả HTTP 200 OK cho Edge",
]
add_bullet_items(slide, Inches(8.5), Inches(2.2), Inches(4.3), Inches(3.5),
                 flow_items, font_size=13, bullet_char="▸", spacing=1.4)

# Storage path
add_rounded_rect(slide, Inches(8.5), Inches(5.0), Inches(4.3), Inches(1.5),
                 fill_color=RGBColor(0xF5, 0xF5, 0xF5), line_color=GRAY)
add_textbox(slide, Inches(8.7), Inches(5.1), Inches(4.0), Inches(0.3),
            "Cấu trúc lưu trữ:", font_size=13, font_color=DARK_BLUE, bold=True)
add_textbox(slide, Inches(8.7), Inches(5.4), Inches(4.0), Inches(1.0),
            "outputs/cloud_uploads/\n  frames/{filename}\n  clips/{filename}\n  events/{filename}",
            font_size=12, font_color=DARK_GRAY, font_name="Consolas", line_spacing=1.3)

add_speaker_notes(slide, """Bảng Alerts là bảng chính trong cơ sở dữ liệu, chứa toàn bộ thông tin về các sự kiện vi phạm.
Mỗi bản ghi bao gồm: id tự tăng, loại vi phạm, thời gian, độ tin cậy, mã thiết bị, đường dẫn ảnh và video bằng chứng, cùng trạng thái xác thực.
Khi Edge gửi dữ liệu lên Cloud, request được gửi dưới dạng multipart/form-data đến POST /alerts.
Backend kiểm tra tính hợp lệ bằng Pydantic, sau đó lưu file bằng chứng vào File System và ghi metadata vào SQLite.
Cấu trúc lưu trữ được phân loại theo frames, clips và events để dễ truy vết.""")


# ════════════════════════════════════════════════════════════════════
# SLIDE 22: Dashboard giám sát
# ════════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_slide_header(slide, "DASHBOARD GIÁM SÁT", "Chương 3.4.2: Giao diện giám sát")

# Dashboard features
features = [
    ("📊 Dashboard Tổng Quan", "Biểu đồ xu hướng cảnh báo 24h\nPhân bố sự kiện theo loại\nTự động làm mới mỗi 5 giây", MEDIUM_BLUE),
    ("🔔 Alerts Center", "Quản lý danh sách cảnh báo\nBộ lọc đa tiêu chí, phân trang\nXác thực hàng loạt", TEAL),
    ("📸 Evidence Modal", "Xem bằng chứng ảnh/clip\nĐánh giá thủ công (manual review)\nKích hoạt xác thực SlowFast", ORANGE),
    ("📱 Devices & Settings", "Quản lý thiết bị Edge\nCấu hình hệ thống\nXuất báo cáo", PURPLE),
]

for i, (title, desc, color) in enumerate(features):
    x = Inches(0.8) + Inches(i * 3.05)
    y = Inches(1.8)

    add_rounded_rect(slide, x, y, Inches(2.8), Inches(2.8), fill_color=WHITE, line_color=color)

    # Header
    add_shape(slide, x, y, Inches(2.8), Inches(0.6), fill_color=color)
    add_textbox(slide, x, y + Inches(0.1), Inches(2.8), Inches(0.4),
                title, font_size=14, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, x + Inches(0.2), y + Inches(0.8), Inches(2.4), Inches(1.8),
                desc, font_size=13, font_color=DARK_GRAY, alignment=PP_ALIGN.CENTER,
                line_spacing=1.4)

# Tech stack
add_textbox(slide, Inches(0.8), Inches(5.0), Inches(12), Inches(0.4),
            "Công nghệ frontend:", font_size=16, font_color=DARK_BLUE, bold=True)

tech_stack = [
    ("React 19", "SPA Framework", MEDIUM_BLUE),
    ("Vite 7", "Build Tool", TEAL),
    ("Recharts", "Biểu đồ Bar/Pie", ORANGE),
    ("Lucide React", "Icon system", PURPLE),
]

for i, (name, desc, color) in enumerate(tech_stack):
    x = Inches(0.8) + Inches(i * 3.05)
    y = Inches(5.5)
    add_rounded_rect(slide, x, y, Inches(2.8), Inches(0.8),
                     fill_color=RGBColor(0xF5, 0xF5, 0xF5), line_color=color)
    add_textbox(slide, x + Inches(0.1), y + Inches(0.05), Inches(2.6), Inches(0.35),
                name, font_size=14, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.1), y + Inches(0.4), Inches(2.6), Inches(0.3),
                desc, font_size=11, font_color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

add_speaker_notes(slide, """Dashboard giám sát được xây dựng bằng React 19 với Vite 7, tạo thành ứng dụng SPA.
Giao diện gồm 4 phần chức năng chính:
Dashboard Tổng Quan hiển thị biểu đồ xu hướng cảnh báo 24h và phân bố sự kiện, tự động làm mới mỗi 5 giây.
Alerts Center cho phép quản lý danh sách cảnh báo với bộ lọc đa tiêu chí, phân trang và xác thực hàng loạt.
Evidence Modal cho phép xem bằng chứng ảnh/clip, đánh giá thủ công và kích hoạt xác thực SlowFast.
Devices & Settings quản lý thiết bị Edge, cấu hình hệ thống và xuất báo cáo.
Thư viện Recharts được sử dụng để vẽ biểu đồ, Lucide React cho hệ thống icon thống nhất.""")


# ════════════════════════════════════════════════════════════════════
# SLIDE 23: Môi trường triển khai và thực nghiệm
# ════════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_slide_header(slide, "MÔI TRƯỜNG TRIỂN KHAI VÀ THỰC NGHIỆM", "Chương 4: Xây dựng và thực nghiệm")

# Edge deployment
add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.4),
            "⚙️ Triển khai phân hệ Edge", font_size=16, font_color=MEDIUM_BLUE, bold=True)

edge_deploy = [
    "Ngôn ngữ: Python, môi trường venv",
    "YOLO: Ultralytics YOLO11m (best.pt)",
    "MediaPipe Pose: 13 landmarks",
    "OpenCV: Xử lý video I/O",
    "CUDA: Tăng tốc GPU (nếu có)",
    "CLI: python -m app.edge.main_edge --source",
    "Chạy headless mode (show_window=false)",
]
add_bullet_items(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(3.5),
                 edge_deploy, font_size=13, bullet_char="▸", spacing=1.3)

# Cloud deployment
add_textbox(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(0.4),
            "☁️ Triển khai phân hệ Cloud", font_size=16, font_color=GREEN, bold=True)

cloud_deploy = [
    "FastAPI: Backend REST API",
    "SQLite + SQLAlchemy: ORM",
    "Pydantic: Schema validation",
    "File System: Lưu bằng chứng",
    "React 19 + Vite 7: Dashboard",
    "config.yaml: Cấu hình runtime",
    "Chạy trên Windows (phát triển)",
]
add_bullet_items(slide, Inches(7), Inches(2.2), Inches(5.5), Inches(3.5),
                 cloud_deploy, font_size=13, bullet_char="▸", spacing=1.3)

# Dataset info
add_textbox(slide, Inches(0.8), Inches(5.0), Inches(12), Inches(0.4),
            "Dữ liệu huấn luyện:", font_size=16, font_color=DARK_BLUE, bold=True)

dataset_info = [
    ("Kaggle State Farm", "~102K frames", MEDIUM_BLUE),
    ("Roboflow Universe", "Multi-source", TEAL),
    ("Custom Local Data", "Bổ sung đặc thù", ORANGE),
]

for i, (source, desc, color) in enumerate(dataset_info):
    x = Inches(0.8) + Inches(i * 4.1)
    y = Inches(5.5)
    add_rounded_rect(slide, x, y, Inches(3.8), Inches(0.8),
                     fill_color=RGBColor(0xF5, 0xF5, 0xF5), line_color=color)
    add_textbox(slide, x + Inches(0.1), y + Inches(0.05), Inches(3.6), Inches(0.35),
                source, font_size=14, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.1), y + Inches(0.4), Inches(3.6), Inches(0.3),
                desc, font_size=12, font_color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# 4 classes
add_textbox(slide, Inches(0.8), Inches(6.5), Inches(12), Inches(0.4),
            "4 lớp đối tượng: phone  |  smoking  |  seatbelt  |  no-seatbelt",
            font_size=14, font_color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

add_speaker_notes(slide, """Về môi trường triển khai, phân hệ Edge sử dụng Python với các thư viện Ultralytics YOLO, MediaPipe Pose và OpenCV.
Hệ thống hỗ trợ CUDA để tăng tốc GPU và có chế độ headless mode cho triển khai thực tế.
Phân hệ Cloud sử dụng FastAPI, SQLite kết hợp SQLAlchemy, và React 19 cho Dashboard.
Dữ liệu huấn luyện được kết hợp từ 3 nguồn: Kaggle State Farm với khoảng 102K frames, Roboflow Universe, và dữ liệu tự thu thập.
Mô hình được huấn luyện với 4 lớp đối tượng: phone, smoking, seatbelt và no-seatbelt.""")


# ════════════════════════════════════════════════════════════════════
# SLIDE 24: Kết quả đánh giá hiệu năng FPS
# ════════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_slide_header(slide, "KẾT QUẢ ĐÁNH GIÁ HIỆU NĂNG FPS", "Chương 4.5.1: Bảng 4.1")

# FPS table
fps_data = [
    ["Kịch bản", "YOLO", "Pose", "Frame\nSkipping", "Kích thước", "FPS"],
    ["Cấu hình 1 (Cơ sở)", "Bật", "Tắt", "Không", "Gốc (848×480)", "2.34"],
    ["Cấu hình 2 (Toàn tải)", "Bật", "Bật", "Không", "Gốc (848×480)", "2.20"],
    ["Cấu hình 3 (Tối ưu chu kỳ)", "Bật", "Bật", "Có (N=3)", "Gốc (848×480)", "3.04"],
    ["Cấu hình 4 (Tối ưu toàn diện)", "Bật", "Bật", "Có (N=3)", "Resize (640)", "4.75"],
]

add_table(slide, Inches(0.5), Inches(1.6), Inches(12), Inches(3.0),
          len(fps_data), 6, fps_data,
          col_widths=[Inches(2.8), Inches(1.2), Inches(1.2), Inches(1.8), Inches(2.5), Inches(1.5)])

# Bar chart simulation
add_textbox(slide, Inches(0.8), Inches(5.0), Inches(5), Inches(0.4),
            "Biểu đồ FPS theo cấu hình:", font_size=16, font_color=DARK_BLUE, bold=True)

fps_values = [2.34, 2.20, 3.04, 4.75]
fps_labels = ["C1\nCơ sở", "C2\nToàn tải", "C3\nTối ưu\chu kỳ", "C4\nTối ưu\ntoàn diện"]
bar_colors = [GRAY, RED, ORANGE, GREEN]

max_fps = 5.0
for i, (val, label, color) in enumerate(zip(fps_values, fps_labels, bar_colors)):
    x = Inches(1.0) + Inches(i * 1.5)
    bar_height = Inches(1.5) * (val / max_fps)
    bar_y = Inches(6.6) - bar_height

    # Bar
    add_shape(slide, x, bar_y, Inches(1.0), bar_height, fill_color=color)
    # Value
    add_textbox(slide, x, bar_y - Inches(0.3), Inches(1.0), Inches(0.3),
                f"{val} FPS", font_size=12, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)
    # Label
    add_textbox(slide, x - Inches(0.2), Inches(6.65), Inches(1.4), Inches(0.5),
                label, font_size=10, font_color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# Key finding
add_rounded_rect(slide, Inches(7), Inches(5.0), Inches(5.5), Inches(1.5),
                 fill_color=RGBColor(0xE8, 0xF5, 0xE9), line_color=GREEN)
add_textbox(slide, Inches(7.3), Inches(5.1), Inches(5.0), Inches(0.4),
            "📈 Kết quả chính:", font_size=16, font_color=GREEN, bold=True)
add_textbox(slide, Inches(7.3), Inches(5.5), Inches(5.0), Inches(0.9),
            "• FPS tăng hơn 2 lần nhờ tổ hợp frame skipping + resize\n"
            "• Từ 2.34 FPS (cơ sở) lên 4.75 FPS (tối ưu toàn diện)\n"
            "• Benchmark chạy trên CPU, trên Edge (Jetson) dự kiến cao hơn",
            font_size=13, font_color=DARK_GRAY, line_spacing=1.4)

add_speaker_notes(slide, """Bảng 4.1 cho thấy kết quả đánh giá hiệu năng FPS theo 4 cấu hình khác nhau.
Cấu hình 1 chỉ chạy YOLO, đạt 2.34 FPS. Khi thêm MediaPipe Pose, FPS giảm nhẹ xuống 2.20 do tải tính toán tăng.
Áp dụng frame skipping với N=3, FPS cải thiện lên 3.04. Kết hợp thêm resize, FPS đạt 4.75 - tăng hơn 2 lần so với cấu hình cơ sở.
Kết quả chứng minh hiệu quả của tổ hợp chiến lược tối ưu frame skipping và resize trong việc cân bằng giữa độ chính xác và hiệu năng.
Lưu ý rằng benchmark này chạy trên máy tính CPU, trên thiết bị Edge như Jetson, FPS dự kiến sẽ cao hơn.""")


# ════════════════════════════════════════════════════════════════════
# SLIDE 25: Kết quả đánh giá độ chính xác
# ════════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_slide_header(slide, "KẾT QUẢ ĐÁNH GIÁ ĐỘ CHÍNH XÁC", "Chương 4.5.2: Bảng 4.2 & 4.3")

# YOLO metrics
add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.4),
            "Đánh giá mô hình YOLO (Bảng 4.2)", font_size=16, font_color=DARK_BLUE, bold=True)

yolo_metrics = [
    ["Chỉ số", "Giá trị", "Epoch"],
    ["Precision", "87.89%", "99"],
    ["Recall", "82.90%", "55"],
    ["mAP50", "88.93%", "100"],
    ["mAP50-95", "59.81%", "100"],
]

add_table(slide, Inches(0.5), Inches(2.1), Inches(5.5), Inches(2.5),
          len(yolo_metrics), 3, yolo_metrics,
          col_widths=[Inches(2.0), Inches(1.8), Inches(1.7)])

# Behavior detection metrics
add_textbox(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(0.4),
            "Hiệu suất nhận diện hành vi (Bảng 4.3)", font_size=16, font_color=DARK_BLUE, bold=True)

behavior_data = [
    ["Hành vi", "TP", "FP", "FN", "Precision", "Recall", "F1"],
    ["Điện thoại", "1165", "90", "88", "92.86%", "92.98%", "92.92%"],
    ["Hút thuốc", "523", "102", "87", "83.69%", "85.74%", "84.70%"],
    ["Ko dây AT", "600", "47", "109", "92.73%", "84.59%", "88.47%"],
]

add_table(slide, Inches(6.5), Inches(2.1), Inches(6.3), Inches(2.5),
          len(behavior_data), 7, behavior_data,
          col_widths=[Inches(1.2), Inches(0.8), Inches(0.7), Inches(0.7), Inches(1.0), Inches(1.0), Inches(0.9)])

# F1 Score bar chart simulation
add_textbox(slide, Inches(0.8), Inches(5.0), Inches(12), Inches(0.4),
            "F1-Score theo từng hành vi:", font_size=16, font_color=DARK_BLUE, bold=True)

f1_data = [
    ("Sử dụng\nđiện thoại", 92.92, GREEN),
    ("Hút thuốc", 84.70, ORANGE),
    ("Ko thắt dây\nan toàn", 88.47, MEDIUM_BLUE),
]

for i, (label, value, color) in enumerate(f1_data):
    x = Inches(1.5) + Inches(i * 3.5)
    bar_height = Inches(1.3) * (value / 100)
    bar_y = Inches(6.5) - bar_height

    add_shape(slide, x, bar_y, Inches(2.5), bar_height, fill_color=color)
    add_textbox(slide, x, bar_y - Inches(0.3), Inches(2.5), Inches(0.3),
                f"{value}%", font_size=16, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, Inches(6.55), Inches(2.5), Inches(0.5),
                label, font_size=12, font_color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

add_speaker_notes(slide, """Về đánh giá mô hình YOLO, chỉ số Precision đạt 87.89%, Recall 82.90%, mAP50 đạt 88.93%.
Tuy nhiên, mAP50-95 chỉ đạt 59.81%, phản ánh độ chính xác định vị bounding box còn hạn chế ở các ngưỡng IoU nghiêm ngặt.
Về hiệu suất nhận diện hành vi, hành vi sử dụng điện thoại đạt F1-Score cao nhất 92.92%.
Hành vi hút thuốc có F1-Score thấp nhất 84.70% do đối tượng nhỏ và dễ bị che khuất.
Hành vi không thắt dây an toàn đạt F1-Score 88.47%.
Nhìn chung, hệ thống đạt hiệu suất nhận diện tương đối tốt trên tập kiểm thử.""")


# ════════════════════════════════════════════════════════════════════
# SLIDE 26: Phân tích lỗi và đánh giá
# ════════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_slide_header(slide, "PHÂN TÍCH LỖI VÀ ĐÁNH GIÁ", "Chương 4.5.3 & 4.5.4")

# Error types
error_types = [
    ("Nhiễu thị giác", "⚠", RED, [
        "Vật thể hình hộp chữ nhật màu tối\n(ví, sạc dự phòng) bị nhận nhầm là điện thoại",
    ]),
    ("Suy thoái môi trường", "⚡", ORANGE, [
        "Ngược sáng gắt hoặc thiếu sáng khiến\nMediaPipe mất landmarks",
        "Dây an toàn hòa lẫn với áo khoác",
    ]),
    ("Biến dạng phối cảnh", "🔄", PURPLE, [
        "Camera bị xô lệch do xe xóc\nlàm Driver ROI dịch chuyển sai lệch",
    ]),
]

for i, (title, icon, color, items) in enumerate(error_types):
    x = Inches(0.8) + Inches(i * 4.1)
    y = Inches(1.8)

    add_rounded_rect(slide, x, y, Inches(3.8), Inches(2.5), fill_color=WHITE, line_color=color)
    add_shape(slide, x, y, Inches(3.8), Inches(0.5), fill_color=color)
    add_textbox(slide, x + Inches(0.1), y + Inches(0.05), Inches(3.6), Inches(0.4),
                f"{icon} {title}", font_size=15, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    for j, item in enumerate(items):
        add_textbox(slide, x + Inches(0.3), y + Inches(0.7 + j * 0.8), Inches(3.2), Inches(0.7),
                    item, font_size=12, font_color=DARK_GRAY, line_spacing=1.3)

# Evidence mechanism
add_textbox(slide, Inches(0.8), Inches(4.7), Inches(12), Inches(0.4),
            "Cơ chế lưu trữ và truy vết bằng chứng:", font_size=16, font_color=DARK_BLUE, bold=True)

evidence_items = [
    ("Tính pháp lý", "Đối soát thủ công trước khi xử phạt tài xế", MEDIUM_BLUE),
    ("Vòng lặp dữ liệu", "Ảnh FP/FN được đẩy về kho dữ liệu để tái huấn luyện", TEAL),
    ("Tinh chỉnh tham số", "Phân tích bằng chứng để điều chỉnh ngưỡng và rules", ORANGE),
]

for i, (title, desc, color) in enumerate(evidence_items):
    x = Inches(0.8) + Inches(i * 4.1)
    y = Inches(5.3)
    add_rounded_rect(slide, x, y, Inches(3.8), Inches(1.2),
                     fill_color=RGBColor(0xE3, 0xF2, 0xFD), line_color=color)
    add_textbox(slide, x + Inches(0.1), y + Inches(0.1), Inches(3.6), Inches(0.35),
                title, font_size=14, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.1), y + Inches(0.5), Inches(3.6), Inches(0.6),
                desc, font_size=12, font_color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

add_speaker_notes(slide, """Trong quá trình thực nghiệm, em ghi nhận 3 nhóm lỗi chính.
Thứ nhất là nhiễu thị giác, khi các vật thể hình hộp chữ nhật màu tối như ví, sạc dự phòng bị nhận nhầm là điện thoại.
Thứ hai là suy thoái môi trường, khi ngược sáng hoặc thiếu sáng khiến MediaPipe mất landmarks, hoặc dây an toàn hòa lẫn với áo.
Thứ ba là biến dạng phối cảnh, khi camera bị xô lệch do xe xóc làm Driver ROI dịch chuyển sai lệch.
Tuy nhiên, cơ chế lưu trữ bằng chứng giúp hệ thống có thể truy vết và cải thiện liên tục.
Ảnh FP/FN được đẩy về kho dữ liệu để tái huấn luyện, và phân tích bằng chứng giúp điều chỉnh ngưỡng và rules.""")


# ════════════════════════════════════════════════════════════════════
# SLIDE 27: Kết luận
# ════════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_slide_header(slide, "KẾT LUẬN", "Chương 5.1: Kết quả đạt được")

conclusions = [
    ("🏗️ Thiết kế kiến trúc", "Hệ thống DMS theo kiến trúc Hybrid Edge–Cloud\nvới phân vai rõ ràng giữa Edge và Cloud", MEDIUM_BLUE),
    ("⚡ Xử lý tại Edge", "Kết hợp YOLO + MediaPipe Pose\ncùng Behavior Rules Engine\ncảnh báo gần thời gian thực", TEAL),
    ("☁️ Quản lý trên Cloud", "FastAPI + SQLite + Dashboard\nlưu trữ, quản lý và giám sát\ncảnh báo tập trung", GREEN),
    ("✅ Tính khả thi", "Chứng minh qua thực nghiệm:\nFPS cải thiện, F1-Score cao\nCơ chế giảm cảnh báo sai hiệu quả", ORANGE),
]

for i, (title, desc, color) in enumerate(conclusions):
    x = Inches(0.8) + Inches(i * 3.05)
    y = Inches(1.8)

    add_rounded_rect(slide, x, y, Inches(2.8), Inches(3.5), fill_color=WHITE, line_color=color)

    add_shape(slide, x, y, Inches(2.8), Inches(0.6), fill_color=color)
    add_textbox(slide, x, y + Inches(0.1), Inches(2.8), Inches(0.4),
                title, font_size=14, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, x + Inches(0.2), y + Inches(0.8), Inches(2.4), Inches(2.5),
                desc, font_size=14, font_color=DARK_GRAY, alignment=PP_ALIGN.CENTER,
                line_spacing=1.4)

# Summary box
add_rounded_rect(slide, Inches(0.8), Inches(5.6), Inches(11.5), Inches(1.2),
                 fill_color=RGBColor(0xE3, 0xF2, 0xFD), line_color=MEDIUM_BLUE)
add_textbox(slide, Inches(1.2), Inches(5.7), Inches(11), Inches(1.0),
            "Đề tài đã hoàn thành các mục tiêu chính ở mức nguyên mẫu thử nghiệm. Hệ thống chứng minh "
            "tính khả thi của kiến trúc Hybrid Edge–Cloud trong bài toán nhận diện hành vi tài xế xe khách, "
            "vừa đáp ứng yêu cầu phản hồi thời gian thực tại Edge, vừa hỗ trợ quản trị tập trung trên Cloud.",
            font_size=15, font_color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER, line_spacing=1.4)

add_speaker_notes(slide, """Tóm lại, đề tài đã đạt được 4 kết quả chính.
Thứ nhất, đã thiết kế thành công kiến trúc Hybrid Edge-Cloud với phân vai rõ ràng giữa Edge xử lý thời gian thực và Cloud quản lý tập trung.
Thứ hai, phân hệ Edge kết hợp YOLO và MediaPipe Pose cùng Behavior Rules Engine, có khả năng cảnh báo gần thời gian thực.
Thứ ba, phân hệ Cloud xây dựng bằng FastAPI, SQLite và Dashboard, hỗ trợ lưu trữ và giám sát cảnh báo.
Thứ tư, kết quả thực nghiệm chứng minh tính khả thi với FPS cải thiện và F1-Score cao cho các hành vi.
Nhìn chung, đề tài đã hoàn thành các mục tiêu chính ở mức nguyên mẫu thử nghiệm.""")


# ════════════════════════════════════════════════════════════════════
# SLIDE 28: Hạn chế và hướng phát triển
# ════════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_slide_header(slide, "HẠN CHẾ VÀ HƯỚNG PHÁT TRIỂN", "Chương 5.2 & 5.3")

# Left - Limitations
add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.4),
            "⚠️ Hạn chế hiện tại", font_size=18, font_color=RED, bold=True)

limitations = [
    "Nhạy cảm với môi trường quang học (ngược sáng, thiếu sáng)",
    "Dataset chưa đủ độ phủ cho mọi tình huống thực tế",
    "Luật tất định với ngưỡng thủ công, thiếu thích ứng động",
    "Chỉ giới hạn 3 hành vi, chưa mở rộng micro-biometric",
    "SlowFast mới ở baseline, chưa fine-tune cho hành vi tài xế",
    "Chưa kiểm thử trên bo mạch nhúng (Jetson) dài hạn",
]
add_bullet_items(slide, Inches(0.8), Inches(2.2), Inches(5.8), Inches(4.5),
                 limitations, font_size=13, bullet_char="✗", spacing=1.4, font_color=RED)

# Right - Future directions
add_textbox(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(0.4),
            "🚀 Hướng phát triển", font_size=18, font_color=GREEN, bold=True)

future = [
    "Mở rộng dataset đa miền, tích hợp Camera hồng ngoại (IR)",
    "Thay thế bằng RT-DETR, Vision Transformers (ViT)",
    "Mở rộng: Gaze Estimation, Head Pose, PERCLOS (buồn ngủ)",
    "Triển khai trên Jetson Orin, tối ưu TensorRT INT8",
    "Hoàn thiện SlowFast pipeline, fine-tune chuyên sâu",
    "Triển khai thực địa trên đội xe doanh nghiệp",
]
add_bullet_items(slide, Inches(7), Inches(2.2), Inches(5.8), Inches(4.5),
                 future, font_size=13, bullet_char="✓", spacing=1.4, font_color=GREEN)

# Bottom summary
add_rounded_rect(slide, Inches(1.5), Inches(6.0), Inches(10), Inches(0.7),
                 fill_color=RGBColor(0xFF, 0xF3, 0xE0), line_color=ORANGE)
add_textbox(slide, Inches(1.8), Inches(6.1), Inches(9.5), Inches(0.5),
            "Đề tài tạo nền tảng để tiếp tục cải thiện độ chính xác, khả năng tối ưu và mức độ ổn định khi triển khai thực tế",
            font_size=14, font_color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

add_speaker_notes(slide, """Về hạn chế, hệ thống còn nhạy cảm với môi trường quang học, dataset chưa đủ độ phủ.
Luật tất định với ngưỡng thủ công thiếu tính thích ứng động. Hiện tại chỉ giới hạn 3 hành vi.
SlowFast mới ở mức baseline, chưa fine-tune chuyên sâu. Chưa kiểm thử trên bo mạch nhúng dài hạn.
Về hướng phát triển, thứ nhất là mở rộng dataset và tích hợp Camera hồng ngoại cho giám sát 24/7.
Thứ hai là thay thế bằng các kiến trúc hiện đại hơn như RT-DETR hoặc Vision Transformers.
Thứ ba là mở rộng sang các hành vi sinh trắc học như Gaze Estimation, Head Pose và PERCLOS cho phát hiện buồn ngủ.
Thứ tư là triển khai trên Jetson Orin với tối ưu TensorRT INT8.
Thứ năm là hoàn thiện SlowFast pipeline. Cuối cùng là triển khai thực địa trên đội xe doanh nghiệp.""")


# ════════════════════════════════════════════════════════════════════
# SLIDE 29: Cảm ơn
# ════════════════════════════════════════════════════════════════════
slide = add_blank_slide()

# Full background
add_shape(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, fill_color=DARK_BLUE)

# Decorative elements
add_shape(slide, 0, 0, SLIDE_WIDTH, Inches(0.08), fill_color=LIGHT_BLUE)
add_shape(slide, 0, Inches(7.2), SLIDE_WIDTH, Inches(0.3), fill_color=LIGHT_BLUE)

# Thank you text
add_textbox(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1.0),
            "EM XIN CHÂN THÀNH CẢM ƠN", font_size=40, font_color=WHITE,
            bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(0.8),
            "quý thầy cô trong hội đồng đã lắng nghe và góp ý",
            font_size=24, font_color=ACCENT_BLUE,
            alignment=PP_ALIGN.CENTER)

# Separator
add_shape(slide, Inches(5), Inches(3.8), Inches(3), Inches(0.04), fill_color=LIGHT_BLUE)

# Q&A
add_textbox(slide, Inches(1.5), Inches(4.2), Inches(10), Inches(0.8),
            "Q & A", font_size=48, font_color=ORANGE,
            bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1.5), Inches(5.2), Inches(10), Inches(0.6),
            "Em sẵn sàng giải đáp mọi thắc mắc của quý thầy cô",
            font_size=18, font_color=RGBColor(0xBB, 0xDE, 0xFB),
            alignment=PP_ALIGN.CENTER)

# Info
add_textbox(slide, Inches(1.5), Inches(6.0), Inches(10), Inches(0.8),
            "Sinh viên: La Đại Lộc  |  MSSV: 4551050116\nGVHD: TS. Nguyễn Thanh Bình  |  Trường Đại học Quy Nhơn",
            font_size=16, font_color=RGBColor(0x90, 0xCA, 0xF9),
            alignment=PP_ALIGN.CENTER, line_spacing=1.5)

add_speaker_notes(slide, """Em xin chân thành cảm ơn quý thầy cô trong hội đồng đã dành thời gian lắng nghe phần trình bày của em.
Em rất mong nhận được sự góp ý và nhận xét của quý thầy cô để đề tài được hoàn thiện hơn.
Nếu quý thầy cô có câu hỏi nào, em sẵn sàng giải đáp ạ.
Em xin cảm ơn!""")


# ════════════════════════════════════════════════════════════════════
# Save
# ════════════════════════════════════════════════════════════════════
output_path = r"D:\QNU\KLTN\driver_behavior_hybrid\backend\baocao\slide_khoa_luan_driver_behavior_hybrid.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
